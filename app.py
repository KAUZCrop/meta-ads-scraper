# ============================================================
# AD INTEL v4.0 — Meta Ad Library Intelligence Board
# ============================================================
import sys, asyncio, subprocess, time, uuid, io, json, requests, base64, sqlite3, random
from concurrent.futures import ThreadPoolExecutor, as_completed
import streamlit as st
from urllib.parse import quote, urlparse
from playwright.sync_api import sync_playwright

if sys.platform.startswith("win"):
    asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())

@st.cache_resource
def ensure_browser():
    subprocess.run([sys.executable, "-m", "playwright", "install", "chromium"], check=True)
    return True

def get_api_key():
    try:
        return st.secrets["ANTHROPIC_API_KEY"]
    except Exception:
        return ""

API_KEY = get_api_key()
DB_PATH = "adintel.db"

st.set_page_config(page_title="AD INTEL", page_icon="◼", layout="wide")


# ============================================================
# 로딩 스플래시 스크린
# ============================================================
st.markdown("""
<div id="adintel-loader">
  <div class="al-bg-grid"></div>
  <div class="al-inner">
    <div class="al-logo-wrap">
      <div class="al-logo">AD<span>INTEL</span></div>
      <div class="al-ver-badge">v4.0</div>
    </div>
    <div class="al-tagline">META AD LIBRARY INTELLIGENCE BOARD</div>
    <div class="al-divider"></div>
    <div class="al-bar-wrap">
      <div class="al-bar"></div>
    </div>
    <div class="al-dots">
      <span></span><span></span><span></span>
    </div>
  </div>
</div>

<style>
/* ── 전체 오버레이 ── */
#adintel-loader {
  position: fixed;
  inset: 0;
  z-index: 999999;
  background: #0b0d11;
  display: flex;
  align-items: center;
  justify-content: center;
  animation: al-fadeout 0.55s cubic-bezier(.4,0,.2,1) 2.6s forwards;
  pointer-events: none;
}

/* ── 격자 배경 ── */
.al-bg-grid {
  position: absolute;
  inset: 0;
  background-image:
    linear-gradient(rgba(26,109,255,.06) 1px, transparent 1px),
    linear-gradient(90deg, rgba(26,109,255,.06) 1px, transparent 1px);
  background-size: 48px 48px;
  animation: al-grid-in 1.2s ease forwards;
  opacity: 0;
}

/* ── 중앙 콘텐츠 ── */
.al-inner {
  display: flex;
  flex-direction: column;
  align-items: center;
  gap: 0;
  animation: al-rise 0.7s cubic-bezier(.2,.8,.4,1) 0.1s forwards;
  opacity: 0;
  transform: translateY(18px);
}

/* ── 로고 행 ── */
.al-logo-wrap {
  display: flex;
  align-items: center;
  gap: 12px;
  margin-bottom: 14px;
}

.al-logo {
  font-family: 'Pretendard', -apple-system, sans-serif;
  font-size: 52px;
  font-weight: 900;
  letter-spacing: -1.5px;
  color: #e8eaf0;
  line-height: 1;
}

.al-logo span {
  color: #4f8aff;
}

.al-ver-badge {
  font-family: 'Pretendard', -apple-system, sans-serif;
  font-size: 11px;
  font-weight: 600;
  color: #4f8aff;
  border: 1.5px solid rgba(79,138,255,.35);
  border-radius: 6px;
  padding: 3px 9px;
  letter-spacing: 0.5px;
  background: rgba(79,138,255,.08);
  align-self: flex-start;
  margin-top: 4px;
}

/* ── 태그라인 ── */
.al-tagline {
  font-family: 'Pretendard', -apple-system, sans-serif;
  font-size: 10px;
  font-weight: 500;
  letter-spacing: 3px;
  color: #4b5263;
  text-transform: uppercase;
  margin-bottom: 28px;
}

/* ── 가로 구분선 ── */
.al-divider {
  width: 280px;
  height: 1px;
  background: linear-gradient(90deg, transparent, rgba(79,138,255,.5), transparent);
  margin-bottom: 24px;
}

/* ── 프로그레스 바 ── */
.al-bar-wrap {
  width: 280px;
  height: 2px;
  background: rgba(255,255,255,.06);
  border-radius: 2px;
  overflow: hidden;
  margin-bottom: 20px;
}

.al-bar {
  height: 100%;
  width: 0%;
  background: linear-gradient(90deg, #1a6dff, #5b9dff);
  border-radius: 2px;
  animation: al-progress 2.4s cubic-bezier(.1,.6,.4,1) 0.4s forwards;
  box-shadow: 0 0 8px rgba(79,138,255,.6);
}

/* ── 점 세 개 인디케이터 ── */
.al-dots {
  display: flex;
  gap: 6px;
}

.al-dots span {
  width: 4px;
  height: 4px;
  border-radius: 50%;
  background: #4f8aff;
  opacity: 0.3;
  animation: al-pulse 1.2s ease-in-out infinite;
}

.al-dots span:nth-child(1) { animation-delay: 0s; }
.al-dots span:nth-child(2) { animation-delay: 0.2s; }
.al-dots span:nth-child(3) { animation-delay: 0.4s; }

/* ── 키프레임 ── */
@keyframes al-fadeout {
  0%   { opacity: 1; }
  100% { opacity: 0; pointer-events: none; }
}

@keyframes al-grid-in {
  0%   { opacity: 0; }
  100% { opacity: 1; }
}

@keyframes al-rise {
  0%   { opacity: 0; transform: translateY(18px); }
  100% { opacity: 1; transform: translateY(0); }
}

@keyframes al-progress {
  0%   { width: 0%; }
  30%  { width: 45%; }
  65%  { width: 72%; }
  85%  { width: 88%; }
  100% { width: 100%; }
}

@keyframes al-pulse {
  0%, 100% { opacity: 0.2; transform: scale(1); }
  50%       { opacity: 1;   transform: scale(1.4); }
}
</style>

<script>
// CSS 애니메이션이 끝난 뒤 DOM에서 완전히 제거
(function() {
  const hide = () => {
    const el = document.getElementById('adintel-loader');
    if (el) el.remove();
  };
  setTimeout(hide, 3200);
  // 사용자가 클릭하면 즉시 제거
  document.addEventListener('click', function once() {
    hide();
    document.removeEventListener('click', once);
  }, { once: true });
})();
</script>
""", unsafe_allow_html=True)


# ============================================================
# DB 레이어 — SQLite 영속성
# ============================================================
def db_init():
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.executescript("""
    CREATE TABLE IF NOT EXISTS assets (
        id TEXT PRIMARY KEY,
        keyword TEXT, country TEXT, asset_type TEXT,
        image_url TEXT, source_url TEXT, caption TEXT,
        width INTEGER, height INTEGER, created_at TEXT,
        starred INTEGER DEFAULT 0,
        capture_source TEXT DEFAULT '',
        img_b64 TEXT DEFAULT '',
        ai_json TEXT DEFAULT ''
    );
    CREATE TABLE IF NOT EXISTS hidden (id TEXT PRIMARY KEY);
    CREATE TABLE IF NOT EXISTS history (keyword TEXT PRIMARY KEY, added_at TEXT);
    CREATE TABLE IF NOT EXISTS summary_store (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        summary_json TEXT,
        created_at TEXT
    );
    """)
    con.commit()
    con.close()

def db_upsert_assets(assets: list):
    if not assets:
        return
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    for a in assets:
        cur.execute("""
        INSERT OR REPLACE INTO assets
        (id,keyword,country,asset_type,image_url,source_url,caption,
         width,height,created_at,starred,capture_source,img_b64,ai_json)
        VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?)
        """, (
            a["id"], a["keyword"], a["country"], a["asset_type"],
            a["image_url"], a["source_url"], a["caption"],
            a["width"], a["height"], a["created_at"],
            1 if a.get("starred") else 0,
            a.get("capture_source", ""),
            a.get("img_b64", ""),
            json.dumps(a.get("ai"), ensure_ascii=False) if a.get("ai") else ""
        ))
    con.commit()
    con.close()

def db_load_assets() -> list:
    """img_b64는 세션 메모리 절약을 위해 로드하지 않음. 분석 시 DB에서 직접 조회."""
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.execute("""
        SELECT id,keyword,country,asset_type,image_url,source_url,caption,
               width,height,created_at,starred,capture_source,ai_json
        FROM assets ORDER BY created_at DESC
    """)
    rows = cur.fetchall()
    con.close()
    out = []
    for r in rows:
        ai = None
        if r[12]:
            try:
                ai = json.loads(r[12])
            except Exception:
                pass
        out.append({
            "id": r[0], "keyword": r[1], "country": r[2],
            "asset_type": r[3], "image_url": r[4], "source_url": r[5],
            "caption": r[6], "width": r[7], "height": r[8],
            "created_at": r[9], "starred": bool(r[10]),
            "capture_source": r[11],
            "img_b64": "",   # 세션에는 비워둠
            "ai": ai
        })
    return out

def db_get_img_b64(aid: str) -> str:
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.execute("SELECT img_b64 FROM assets WHERE id=?", (aid,))
    row = cur.fetchone()
    con.close()
    return (row[0] or "") if row else ""

def db_set_field(aid: str, **kwargs):
    """allowed 필드만 업데이트"""
    allowed = {"starred", "img_b64", "ai_json", "capture_source"}
    pairs = [(k, v) for k, v in kwargs.items() if k in allowed]
    if not pairs:
        return
    sql = ", ".join(f"{k}=?" for k, _ in pairs)
    vals = [v for _, v in pairs] + [aid]
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.execute(f"UPDATE assets SET {sql} WHERE id=?", vals)
    con.commit()
    con.close()

def db_add_hidden(aid: str):
    con = sqlite3.connect(DB_PATH)
    con.execute("INSERT OR IGNORE INTO hidden VALUES (?)", (aid,))
    con.commit()
    con.close()

def db_load_hidden() -> set:
    con = sqlite3.connect(DB_PATH)
    rows = con.execute("SELECT id FROM hidden").fetchall()
    con.close()
    return {r[0] for r in rows}

def db_add_history(kw: str):
    con = sqlite3.connect(DB_PATH)
    con.execute("INSERT OR IGNORE INTO history VALUES (?,?)",
                (kw, time.strftime("%Y-%m-%d %H:%M:%S")))
    con.commit()
    con.close()

def db_load_history() -> list:
    con = sqlite3.connect(DB_PATH)
    rows = con.execute("SELECT keyword FROM history ORDER BY added_at").fetchall()
    con.close()
    return [r[0] for r in rows]

def db_save_summary(s: dict):
    con = sqlite3.connect(DB_PATH)
    con.execute("DELETE FROM summary_store")
    con.execute("INSERT INTO summary_store (summary_json,created_at) VALUES (?,?)",
                (json.dumps(s, ensure_ascii=False), time.strftime("%Y-%m-%d %H:%M:%S")))
    con.commit()
    con.close()

def db_load_summary() -> dict | None:
    con = sqlite3.connect(DB_PATH)
    row = con.execute("SELECT summary_json FROM summary_store ORDER BY id DESC LIMIT 1").fetchone()
    con.close()
    if row:
        try:
            return json.loads(row[0])
        except Exception:
            return None
    return None

def db_reset():
    con = sqlite3.connect(DB_PATH)
    con.executescript("""
        DELETE FROM assets;
        DELETE FROM hidden;
        DELETE FROM history;
        DELETE FROM summary_store;
    """)
    con.commit()
    con.close()

db_init()


# ============================================================
# 세션 상태 초기화 — DB에서 로드
# ============================================================
if "initialized" not in st.session_state:
    st.session_state.assets   = db_load_assets()
    st.session_state.hidden   = db_load_hidden()
    st.session_state.history  = db_load_history()
    st.session_state.log            = []
    st.session_state.ai_on          = False
    st.session_state.dark           = False
    st.session_state.selected       = set()
    st.session_state.summary        = db_load_summary()
    st.session_state.analysis_model = "haiku"   # haiku | sonnet
    st.session_state.initialized    = True


# ============================================================
# CSS
# ============================================================
D = st.session_state.dark

LIGHT = """
:root {
    --bg:#fff; --bg2:#f7f8fa; --bg3:#f0f2f5;
    --bd:#e2e5ea; --bd2:#cdd1d8;
    --ac:#1a6dff; --ac2:rgba(26,109,255,0.10);
    --tx:#111318; --tx2:#3a3f4a; --mu:#8a909e;
    --ok:#00a86b; --er:#e8284a; --wn:#e07b00;
    --sh:0 1px 3px rgba(0,0,0,.05);
    --sh2:0 4px 20px rgba(26,109,255,.12);
}"""

DARK = """
:root {
    --bg:#0d0f12; --bg2:#13161b; --bg3:#1a1e26;
    --bd:#252a34; --bd2:#2f3542;
    --ac:#4f8aff; --ac2:rgba(79,138,255,0.12);
    --tx:#e8eaf0; --tx2:#b0b8cc; --mu:#6b7280;
    --ok:#3dffa0; --er:#ff4f6b; --wn:#ffb84f;
    --sh:0 1px 4px rgba(0,0,0,.4);
    --sh2:0 4px 20px rgba(79,138,255,.18);
}"""

st.markdown(f"""<style>
@import url('https://cdn.jsdelivr.net/gh/orioncactus/pretendard/dist/web/static/pretendard.css');

{DARK if D else LIGHT}

html,body,[class*="css"]{{font-family:'Pretendard',sans-serif;background:var(--bg)!important;color:var(--tx)!important;}}
.stApp{{background:var(--bg)!important;}}
[data-testid="stHeader"]{{background:var(--bg)!important;border-bottom:1px solid var(--bd);}}
[data-testid="stSidebar"]{{background:var(--bg2)!important;border-right:1px solid var(--bd)!important;}}
[data-testid="stSidebar"] *{{color:var(--tx)!important;}}
.block-container{{padding-top:4.5rem;padding-bottom:3rem;max-width:1600px;}}

.hdr{{margin-bottom:2rem;padding-bottom:1.2rem;border-bottom:2px solid var(--bd);}}
.hdr-row{{display:flex;align-items:center;gap:10px;}}
.logo{{font-family:'Pretendard',sans-serif;font-size:22px;font-weight:800;letter-spacing:-.5px;color:var(--tx);}}
.logo b{{color:var(--ac);}}
.ver{{font-size:10px;color:var(--mu);border:1px solid var(--bd2);padding:2px 8px;border-radius:4px;background:var(--bg3);}}
.sub{{font-size:10px;color:var(--mu);letter-spacing:2px;margin-top:8px;}}

.slbl{{font-size:9px;color:var(--mu);letter-spacing:1.5px;text-transform:uppercase;margin:8px 0 6px;}}

.srch{{background:var(--bg2);border:1.5px solid var(--bd);border-radius:14px;padding:16px 20px;margin-bottom:1.2rem;}}
.srch-lbl{{font-size:9px;color:var(--mu);letter-spacing:2px;text-transform:uppercase;margin-bottom:8px;}}

.sc{{background:var(--bg);border:1.5px solid var(--bd);border-radius:12px;padding:20px 20px 18px;position:relative;overflow:hidden;box-shadow:var(--sh);transition:.2s;}}
.sc:hover{{border-color:var(--ac);box-shadow:var(--sh2);}}
.sc::before{{content:'';position:absolute;top:0;left:0;right:0;height:3px;background:linear-gradient(90deg,var(--ac),#5b9dff);opacity:0;transition:.2s;}}
.sc:hover::before{{opacity:1;}}
.sc-lbl{{font-size:9px;color:var(--mu);letter-spacing:1px;text-transform:uppercase;margin-bottom:10px;}}
.sc-val{{font-family:'Pretendard',sans-serif;font-size:32px;font-weight:800;color:var(--tx);line-height:1;}}
.sc-sub{{font-size:11px;color:var(--mu);margin-top:8px;}}

.tags{{display:flex;flex-wrap:wrap;gap:6px;margin-bottom:1.2rem;}}
.tag{{font-size:11px;background:var(--ac2);border:1px solid var(--ac2);color:var(--ac);padding:4px 10px;border-radius:20px;display:inline-flex;align-items:center;gap:6px;cursor:pointer;}}
.tag-dot{{width:5px;height:5px;border-radius:50%;background:var(--ok);}}

.sec{{display:flex;align-items:center;justify-content:space-between;margin-bottom:.8rem;padding-bottom:10px;border-bottom:1.5px solid var(--bd);}}
.sec-t{{font-family:'Pretendard',sans-serif;font-size:13px;font-weight:800;color:var(--tx);letter-spacing:.5px;text-transform:uppercase;}}
.sec-n{{font-size:11px;color:var(--mu);}}

.card{{border:1.5px solid var(--bd);border-radius:12px;overflow:hidden;margin-bottom:10px;background:var(--bg);box-shadow:var(--sh);transition:.2s;}}
.card:hover{{border-color:var(--ac);box-shadow:var(--sh2);}}
.card-body{{padding:10px 12px;border-top:1px solid var(--bd);background:var(--bg2);}}
.card-kw{{font-size:10px;color:var(--ac);letter-spacing:.8px;text-transform:uppercase;margin-bottom:4px;}}
.card-meta{{font-size:11px;color:var(--mu);margin-top:4px;}}
.card-cap{{font-size:11px;color:var(--tx2);margin:4px 0;opacity:.8;font-style:italic;line-height:1.4;}}
.bdg{{display:inline-block;font-size:9px;padding:2px 7px;border-radius:20px;margin:0 2px 4px 0;font-weight:600;}}
.b-img{{background:var(--ac2);color:var(--ac);border:1px solid var(--ac2);}}
.b-vid{{background:rgba(124,92,255,.12);color:#7c5cff;border:1px solid rgba(124,92,255,.2);}}
.b-sav{{background:rgba(255,184,79,.12);color:var(--wn);border:1px solid rgba(255,184,79,.2);}}
.b-ai {{background:rgba(61,255,160,.10);color:var(--ok);border:1px solid rgba(61,255,160,.2);}}

.ai-wrap{{padding:0 12px 12px;background:var(--bg2);}}
.ai-box{{background:var(--ac2);border:1px solid var(--ac2);border-radius:10px;padding:10px 12px;}}
.ai-head{{font-size:9px;color:var(--ac);letter-spacing:1.5px;text-transform:uppercase;margin-bottom:8px;font-weight:700;}}
.ai-body{{font-size:12px;color:var(--tx2);line-height:1.75;}}
.ai-tag{{display:inline-block;font-size:10px;background:var(--ac2);color:var(--ac);border-radius:4px;padding:1px 6px;margin:2px 2px 0 0;font-weight:600;}}

.summary-box{{background:var(--bg2);border:2px solid var(--ac);border-radius:16px;padding:24px 28px;margin-bottom:1.5rem;}}
.summary-head{{font-family:'Pretendard',sans-serif;font-size:16px;font-weight:800;color:var(--ac);margin-bottom:16px;}}
.summary-grid{{display:grid;grid-template-columns:1fr 1fr;gap:12px;margin-bottom:14px;}}
.summary-item{{background:var(--bg);border:1px solid var(--bd);border-radius:10px;padding:12px 14px;}}
.summary-item-lbl{{font-size:9px;color:var(--mu);letter-spacing:1px;text-transform:uppercase;margin-bottom:6px;font-weight:700;}}
.summary-item-val{{font-size:13px;color:var(--tx);line-height:1.5;}}
.summary-strategy{{background:var(--ac2);border:1px solid var(--ac2);border-radius:10px;padding:14px 16px;margin-bottom:12px;}}
.summary-strategy-lbl{{font-size:9px;color:var(--ac);letter-spacing:1px;text-transform:uppercase;margin-bottom:8px;font-weight:700;}}
.summary-strategy-val{{font-size:13px;color:var(--tx2);line-height:1.65;}}
.summary-tags{{display:flex;flex-wrap:wrap;gap:6px;}}
.summary-tag{{font-size:10px;background:var(--ac2);color:var(--ac);border:1px solid var(--ac2);border-radius:4px;padding:3px 9px;font-weight:600;}}

.sel-bar{{display:flex;align-items:center;background:var(--bg2);border:1.5px solid var(--bd);border-radius:10px;padding:10px 16px;margin-bottom:1rem;}}
.sel-count{{font-family:'Pretendard',sans-serif;font-size:14px;font-weight:800;color:var(--ac);margin-right:6px;}}
.sel-bar-txt{{font-size:11px;color:var(--tx2);}}

.banner{{display:flex;align-items:center;gap:10px;background:var(--ac2);border:1px solid var(--ac2);border-radius:10px;padding:10px 16px;margin-bottom:1rem;}}
.banner-w{{background:rgba(224,123,0,.06);border-color:rgba(224,123,0,.2);}}
.dot{{width:8px;height:8px;border-radius:50%;flex-shrink:0;}}
.dot-on{{background:var(--ok);box-shadow:0 0 6px var(--ok);}}
.banner-txt{{font-size:11px;color:var(--tx2);}}

.err-box{{background:rgba(232,40,74,.06);border:1px solid rgba(232,40,74,.2);border-radius:10px;padding:10px 14px;margin-bottom:8px;font-size:11px;color:var(--er);}}

.empty{{text-align:center;padding:80px 20px;}}
.empty-t{{font-family:'Pretendard',sans-serif;font-size:16px;font-weight:700;color:var(--mu);margin-bottom:8px;}}
.empty-d{{font-size:13px;color:var(--mu);line-height:1.6;}}

.log{{font-size:11px;color:var(--mu);padding:4px 0;border-bottom:1px solid var(--bd);}}
.log .ts{{color:var(--ac);margin-right:8px;}}
.log .ok{{color:var(--ok);}}
.log .er{{color:var(--er);}}

.persist-badge{{display:inline-flex;align-items:center;gap:5px;font-size:10px;color:var(--ok);background:rgba(0,168,107,.08);border:1px solid rgba(0,168,107,.2);border-radius:6px;padding:3px 8px;}}

.stButton>button{{
    font-family:'Pretendard',sans-serif!important;font-size:13px!important;font-weight:600!important;
    border-radius:8px!important;border:1.5px solid var(--bd2)!important;
    background:var(--bg)!important;color:var(--tx2)!important;height:34px!important;
    transition:.15s!important;white-space:nowrap!important;overflow:hidden!important;
    padding:0 8px!important;line-height:34px!important;
}}
.stButton>button:hover{{border-color:var(--ac)!important;color:var(--ac)!important;background:var(--ac2)!important;}}
.stDownloadButton>button{{
    font-family:'Pretendard',sans-serif!important;font-size:12px!important;
    border-radius:8px!important;border:1.5px solid var(--bd2)!important;
    background:var(--bg)!important;color:var(--tx2)!important;white-space:nowrap!important;
}}

.stTextInput input{{background:var(--bg)!important;border:1.5px solid var(--bd2)!important;border-radius:10px!important;color:var(--tx)!important;font-family:'Pretendard',sans-serif!important;}}
.stTextInput input:focus{{border-color:var(--ac)!important;box-shadow:0 0 0 3px var(--ac2)!important;}}
div[data-baseweb="select"]>div{{background:var(--bg)!important;border-color:var(--bd2)!important;border-radius:8px!important;color:var(--tx)!important;}}

.stSuccess{{background:rgba(0,168,107,.08)!important;border:1px solid rgba(0,168,107,.25)!important;border-radius:10px!important;}}
.stWarning{{background:rgba(224,123,0,.08)!important;border:1px solid rgba(224,123,0,.25)!important;border-radius:10px!important;}}
.stInfo   {{background:var(--ac2)!important;border:1px solid var(--ac2)!important;border-radius:10px!important;}}
.stError  {{background:rgba(232,40,74,.08)!important;border:1px solid rgba(232,40,74,.25)!important;border-radius:10px!important;}}

hr{{border-color:var(--bd)!important;}}
.stCaption{{color:var(--mu)!important;font-size:11px!important;}}
a{{color:var(--ac)!important;}}
::-webkit-scrollbar{{width:4px;height:4px;}}
::-webkit-scrollbar-track{{background:var(--bg2);}}
::-webkit-scrollbar-thumb{{background:var(--bd2);border-radius:4px;}}
.stSpinner>div{{border-top-color:var(--ac)!important;}}
[data-testid="stTabs"]{{border-bottom:2px solid var(--bd);}}
</style>""", unsafe_allow_html=True)


# ============================================================
# 스크래퍼
# ============================================================
def valid(w, h):
    if w < 180 or h < 180:
        return False
    r = w / h if h else 0
    return 0.25 <= r <= 3.0

def make_fp(item):
    url = (item.get("image_url") or "").strip()
    p = urlparse(url)
    return (
        f"{p.netloc}{p.path}",
        (item.get("asset_type") or "").strip(),
        (item.get("caption") or "").strip().lower()[:80],
        item.get("width", 0),
        item.get("height", 0),
    )

def _bytes_to_b64(img_bytes: bytes | None) -> str:
    if not img_bytes:
        return ""
    return base64.b64encode(img_bytes).decode("utf-8")

def _close_modal(page):
    try:
        page.keyboard.press("Escape")
        page.wait_for_timeout(300)
    except Exception:
        pass

def _is_valid_screenshot(img_bytes: bytes | None) -> bool:
    if not img_bytes or len(img_bytes) < 3000:
        return False
    try:
        from PIL import Image, ImageStat
        im = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        w, h = im.size
        if w < 120 or h < 120:
            return False
        stat = ImageStat.Stat(im.resize((32, 32)))
        if max(stat.stddev) < 3:
            return False
        return True
    except Exception:
        return True

def _url_key(url: str) -> str:
    p = urlparse((url or "").strip())
    return f"{p.netloc}{p.path}"

def _find_media_handle_by_url(page, target_url: str, asset_type: str = "image"):
    target_key = _url_key(target_url)
    try:
        handle = page.evaluate_handle("""({targetKey, assetType}) => {
            const key = (u) => {
                try { const p = new URL(u); return p.host + p.pathname; }
                catch(e) { return u || ''; }
            };
            const nodes = assetType === 'video_poster'
                ? Array.from(document.querySelectorAll('video'))
                : Array.from(document.querySelectorAll('img'));
            for (const el of nodes) {
                const src = assetType === 'video_poster'
                    ? (el.poster || '')
                    : (el.currentSrc || el.src || '');
                if (src && key(src) === targetKey) return el;
            }
            return null;
        }""", {"targetKey": target_key, "assetType": asset_type})
        return handle.as_element()
    except Exception:
        return None

def _find_ad_card_container(page, target_key: str):
    """
    이미지 URL을 기준으로 광고 카드 전체 컨테이너를 DOM에서 탐색.
    텍스트 카피·CTA·브랜드명까지 포함된 카드 영역을 반환.
    """
    try:
        handle = page.evaluate_handle("""({targetKey}) => {
            const key = (u) => {
                try { const p = new URL(u); return p.host + p.pathname; }
                catch(e) { return u || ''; }
            };
            const imgs = Array.from(document.querySelectorAll('img'));
            for (const img of imgs) {
                const src = img.currentSrc || img.src || '';
                if (!src || key(src) !== targetKey) continue;
                let el = img;
                for (let i = 0; i < 12; i++) {
                    if (!el.parentElement) break;
                    el = el.parentElement;
                    const rect = el.getBoundingClientRect();
                    const style = window.getComputedStyle(el);
                    if (rect.width > 200 && rect.height > 250
                        && style.overflow !== 'scroll'
                        && style.overflow !== 'auto') {
                        const text = (el.innerText || '').trim();
                        if (text.length > 15) return el;
                    }
                }
                return img;
            }
            return null;
        }""", {"targetKey": target_key})
        return handle.as_element()
    except Exception:
        return None


def _capture_media_by_url_b64(page, target_url: str, asset_type: str = "image") -> tuple[str, str]:
    """
    개선된 3단계 캡처 전략 (quality=85):
    1순위 — 광고 카드 컨테이너 전체 (텍스트+이미지+CTA 포함)
    2순위 — 모달 다이얼로그
    3순위 — 이미지 요소 단독
    """
    target_key = _url_key(target_url)

    # ── 1순위: 광고 카드 컨테이너 ──
    try:
        card = _find_ad_card_container(page, target_key)
        if card:
            card.scroll_into_view_if_needed(timeout=7000)
            page.wait_for_timeout(700)
            shot = card.screenshot(type="jpeg", quality=85, timeout=15000)
            if _is_valid_screenshot(shot):
                return _bytes_to_b64(shot), "card_container"
    except Exception:
        pass

    # ── 2순위: 클릭 후 모달 ──
    element = _find_media_handle_by_url(page, target_url, asset_type)
    if element:
        try:
            element.scroll_into_view_if_needed(timeout=7000)
            page.wait_for_timeout(600)
            element.click(timeout=5000, force=True)
            page.wait_for_timeout(1600)
            dialog = page.locator('div[role="dialog"]').last
            if dialog.count() > 0 and dialog.is_visible():
                shot = dialog.screenshot(type="jpeg", quality=85, timeout=15000)
                _close_modal(page)
                if _is_valid_screenshot(shot):
                    return _bytes_to_b64(shot), "modal"
        except Exception:
            _close_modal(page)

        # ── 3순위: 이미지 요소 단독 ──
        try:
            element = _find_media_handle_by_url(page, target_url, asset_type) or element
            element.scroll_into_view_if_needed(timeout=7000)
            page.wait_for_timeout(500)
            shot = element.screenshot(type="jpeg", quality=85, timeout=12000)
            if _is_valid_screenshot(shot):
                return _bytes_to_b64(shot), "element"
        except Exception:
            pass

    return "", "capture_failed"

def scrape(keyword, country, scrolls, limit):
    ensure_browser()
    search_url = (
        "https://www.facebook.com/ads/library/"
        f"?active_status=all&ad_type=all&country={country}&q={quote(keyword)}"
    )
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True, args=[
            "--no-sandbox", "--disable-dev-shm-usage", "--disable-extensions",
            "--disable-plugins", "--disable-background-networking",
            "--no-first-run", "--disable-default-apps",
        ])
        ctx = browser.new_context(
            viewport={"width": 1440, "height": 2000},
            user_agent=(
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0.0.0 Safari/537.36"
            ),
            locale="ko-KR",
        )
        ctx.route("**/*.{woff,woff2,ttf,otf,eot}", lambda r: r.abort())
        page = ctx.new_page()
        page.goto(search_url, wait_until="domcontentloaded", timeout=90000)
        # 랜덤 딜레이로 봇 탐지 회피
        page.wait_for_timeout(3000 + random.randint(500, 1500))

        def card_count():
            return page.evaluate("""() => {
                const sels = [
                    'div[class*="x8gbvx8"]',
                    'div[data-visualcompletion="ignore-dynamic"]',
                    'div._7jyr',
                    'div[class*="xh8yej3"]'
                ];
                for (const s of sels) {
                    const n = document.querySelectorAll(s).length;
                    if (n > 0) return n;
                }
                return document.querySelectorAll('img[src*="fbcdn"]').length;
            }""")

        prev, stalls = 0, 0
        for _ in range(scrolls):
            page.mouse.wheel(0, random.randint(2800, 3600))
            waited = 0
            while waited < 2500:
                page.wait_for_timeout(random.randint(250, 400))
                waited += 350
                cur = card_count()
                if cur > prev:
                    prev = cur
                    stalls = 0
                    break
            else:
                stalls += 1
                if stalls >= 2:
                    break

        raw = page.evaluate("""() => {
            const out = [];
            const imgs = Array.from(document.querySelectorAll('img'));
            imgs.forEach((img, idx) => {
                const src = img.currentSrc || img.src || '';
                if (!src || src.startsWith('data:')) return;
                out.push({
                    type: 'image', src, idx,
                    w: img.naturalWidth  || img.width  || 0,
                    h: img.naturalHeight || img.height || 0,
                    alt: img.alt || ''
                });
            });
            const vids = Array.from(document.querySelectorAll('video'));
            vids.forEach((v, idx) => {
                const poster = v.poster || '';
                if (!poster) return;
                out.push({
                    type: 'video_poster', src: poster, idx,
                    w: v.videoWidth  || v.clientWidth  || 0,
                    h: v.videoHeight || v.clientHeight || 0,
                    alt: ''
                });
            });
            return out;
        }""")

        ctx.close()
        browser.close()

    seen, collected, now = set(), [], time.strftime("%Y-%m-%d %H:%M:%S")
    for r in raw:
        if len(collected) >= limit:
            break
        src = (r.get("src") or "").strip()
        w, h = int(r.get("w") or 0), int(r.get("h") or 0)
        if not src or not valid(w, h):
            continue
        asset_type = r.get("type", "image")
        asset = {
            "id":             str(uuid.uuid4()),
            "keyword":        keyword,
            "country":        country,
            "asset_type":     asset_type,
            "image_url":      src,
            "source_url":     search_url,
            "caption":        (r.get("alt") or "").strip(),
            "width":          w,
            "height":         h,
            "created_at":     now,
            "starred":        False,
            "ai":             None,
            "img_b64":        "",
            "capture_source": "not_captured",
        }
        fp = make_fp(asset)
        if fp in seen:
            continue
        seen.add(fp)
        collected.append(asset)

    return collected


def capture_screenshots_for_items(items, scrolls=8):
    """선택된 소재만 캡처. 완료 후 DB에 img_b64 저장."""
    if not items:
        return items
    ensure_browser()
    groups = {}
    for item in items:
        if item.get("img_b64"):
            continue
        groups.setdefault(item.get("source_url", ""), []).append(item)

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True, args=[
            "--no-sandbox", "--disable-dev-shm-usage",
            "--disable-extensions", "--disable-plugins",
            "--no-first-run", "--disable-default-apps",
        ])
        ctx = browser.new_context(
            viewport={"width": 1440, "height": 2000},
            user_agent=(
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0.0.0 Safari/537.36"
            ),
        )
        ctx.route("**/*.{woff,woff2,ttf,otf,eot}", lambda r: r.abort())

        for source_url, group_items in groups.items():
            if not source_url:
                for item in group_items:
                    item["capture_source"] = "missing_source_url"
                continue
            page = ctx.new_page()
            try:
                page.goto(source_url, wait_until="domcontentloaded", timeout=90000)
                page.wait_for_timeout(3500)
                target_keys = {_url_key(item.get("image_url")) for item in group_items}
                matched = {}
                for _ in range(max(3, int(scrolls))):
                    raw = page.evaluate("""() => {
                        const out = [];
                        Array.from(document.querySelectorAll('img')).forEach((img, idx) => {
                            const src = img.currentSrc || img.src || '';
                            if (!src || src.startsWith('data:')) return;
                            out.push({type:'image', src, idx});
                        });
                        Array.from(document.querySelectorAll('video')).forEach((v, idx) => {
                            const poster = v.poster || '';
                            if (!poster) return;
                            out.push({type:'video_poster', src:poster, idx});
                        });
                        return out;
                    }""")
                    for r in raw:
                        k = _url_key(r.get("src"))
                        if k in target_keys and k not in matched:
                            matched[k] = r
                    if len(matched) >= len(target_keys):
                        break
                    page.mouse.wheel(0, 3200)
                    page.wait_for_timeout(700)

                for item in group_items:
                    k = _url_key(item.get("image_url"))
                    r = matched.get(k)
                    if not r:
                        item["capture_source"] = "selected_asset_not_found"
                        continue
                    asset_type = r.get("type", "image")
                    img_b64, capture_source = _capture_media_by_url_b64(
                        page,
                        item.get("image_url", ""),
                        "image" if asset_type == "image" else "video_poster",
                    )
                    item["img_b64"] = img_b64
                    item["capture_source"] = capture_source
                    # 캡처 완료 즉시 DB에 저장
                    if img_b64:
                        db_set_field(item["id"], img_b64=img_b64, capture_source=capture_source)
            except Exception as ex:
                for item in group_items:
                    if not item.get("img_b64"):
                        item["capture_source"] = f"capture_error: {str(ex)[:80]}"
            finally:
                try:
                    page.close()
                except Exception:
                    pass

        ctx.close()
        browser.close()

    return items


# ============================================================
# AI 분석
# ============================================================
def _anthropic_model(role: str = "analysis") -> str:
    """
    role: analysis → 사용자 선택 모델 / repair·summary → 항상 haiku (비용 절감)
    """
    if role in ("repair", "summary"):
        return "claude-haiku-4-5-20251001"
    try:
        model_choice = st.session_state.get("analysis_model", "haiku")
    except Exception:
        model_choice = "haiku"
    if model_choice == "sonnet":
        return "claude-sonnet-4-5-20250514"
    return "claude-haiku-4-5-20251001"

# ============================================================
# JSON 파싱 — 4단계 복구 파이프라인
# ============================================================
import re as _re

def _extract_json_block(text: str) -> str | None:
    """마크다운 제거 후 가장 바깥 JSON object 추출."""
    if not text:
        return None
    txt = text.strip()
    # 마크다운 코드블록 제거
    txt = _re.sub(r"```json\s*", "", txt)
    txt = _re.sub(r"```\s*", "", txt)
    txt = txt.strip()
    s = txt.find("{")
    e = txt.rfind("}") + 1
    if s == -1 or e <= 0 or e <= s:
        return None
    return txt[s:e]


def _fix_unescaped_newlines(text: str) -> str:
    """JSON string values에서 이스케이프되지 않은 줄바꿈을 공백으로 교체합니다."""
    result = []
    in_string = False
    escape_next = False
    bs = "\\"
    dq = '"'
    nl = "\n"
    cr = "\r"
    for ch in text:
        if escape_next:
            result.append(ch)
            escape_next = False
        elif ch == bs and in_string:
            result.append(ch)
            escape_next = True
        elif ch == dq:
            in_string = not in_string
            result.append(ch)
        elif ch in (nl, cr) and in_string:
            result.append(" ")
        else:
            result.append(ch)
    return "".join(result)


def _fix_common_json_issues(text: str) -> str:
    """후행 쉼표 제거 및 제어 문자 제거."""
    text = "".join(ch for ch in text if ord(ch) >= 32 or ch in ("\t", "\n", "\r"))
    text = _re.sub(r",(?=\s*[}\]])", "", text)
    return text


def _parse_ai_json(txt: str) -> tuple[dict | None, str | None]:
    """
    4단계 복구 파이프라인:
    1. 직접 파싱
    2. 줄바꿈 수정 후 파싱          ← 가장 흔한 실패 원인
    3. 공통 오류 수정 후 파싱
    4. Claude 복구 요청
    """
    raw_json = _extract_json_block(txt)
    if not raw_json:
        return None, f"JSON 블록을 찾을 수 없습니다. 응답 일부: {txt[:180]}"

    # 단계 1: 직접 파싱
    try:
        return json.loads(raw_json), None
    except json.JSONDecodeError:
        pass

    # 단계 2: 줄바꿈 수정
    step2 = _fix_unescaped_newlines(raw_json)
    try:
        return json.loads(step2), None
    except json.JSONDecodeError:
        pass

    # 단계 3: 공통 오류 수정 (줄바꿈 수정 위에 추가 적용)
    step3 = _fix_common_json_issues(step2)
    try:
        return json.loads(step3), None
    except json.JSONDecodeError as ex3:
        last_err = ex3

    # 단계 4: Claude 복구 (가장 정리된 버전 전달)
    repaired = _repair_json_with_claude(step3)
    if repaired is not None:
        return repaired, None

    return None, f"JSON 복구 실패 (4단계 모두 실패): {last_err}. 원문 일부: {raw_json[:300]}"


def _repair_json_with_claude(raw_json: str) -> dict | None:
    """Claude에게 JSON 문법만 복구 요청. 내용 변경 금지."""
    if not API_KEY or not raw_json:
        return None
    try:
        repair_prompt = f"""아래는 JSON 파싱에 실패한 광고 분석 결과입니다.
내용(값)은 절대 바꾸지 말고, json.loads()로 파싱 가능한 순수 JSON 문법으로만 수정하세요.

수정 규칙:
- 마크다운·백틱·설명문 출력 금지
- 문자열 내 줄바꿈(실제 개행)은 공백 한 칸으로 교체
- 문자열 내 큰따옴표(")는 작은따옴표(')로 교체
- 후행 쉼표(,}} 또는 ,]) 제거
- 누락된 쉼표·따옴표·괄호 복구
- key 이름과 값의 의미 절대 변경 금지

수정할 JSON:
{raw_json[:3000]}
"""
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"x-api-key": API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            json={
                "model": _anthropic_model("repair"),
                "max_tokens": 2000,
                "temperature": 0,
                "messages": [{"role": "user", "content": repair_prompt}],
            },
            timeout=40,
        )
        if resp.status_code != 200:
            return None
        fixed_txt = resp.json()["content"][0]["text"].strip()
        fixed_json = _extract_json_block(fixed_txt)
        if not fixed_json:
            return None
        # 복구 결과에도 4단계 중 1~3 적용
        for attempt in [fixed_json,
                        _fix_unescaped_newlines(fixed_json),
                        _fix_common_json_issues(_fix_unescaped_newlines(fixed_json))]:
            try:
                return json.loads(attempt)
            except Exception:
                continue
        return None
    except Exception:
        return None



# ============================================================
# 분석 로딩 애니메이션 -- 스플래시 스크린과 동일 톤
# ============================================================
def _analysis_status_html(step, current, total, detail=""):
    BG      = "#0d0f12"
    BORDER  = "#252a34"
    TX      = "#e8eaf0"
    MU      = "#6b7280"
    AC      = "#4f8aff"
    AC_DIM  = "rgba(79,138,255,.08)"
    AC_MED  = "rgba(79,138,255,.15)"
    AC_BD   = "rgba(79,138,255,.35)"
    AC_GLOW = "rgba(79,138,255,.6)"
    OK      = "#3dffa0"
    OK_DIM  = "rgba(61,255,160,.12)"
    OK_BD   = "rgba(61,255,160,.3)"
    GRID    = "rgba(26,109,255,.055)"
    TRACK   = "rgba(255,255,255,.06)"

    step_order = {"capture": 0, "analyze": 2, "done": 3}
    step_idx   = step_order.get(step, 0)
    pct        = int((current / max(total, 1)) * 100)
    is_done    = step == "done"

    step_defs = [
        ("01", "CAPTURE",  "Playwright"),
        ("02", "OCR",      "Claude Vision"),
        ("03", "ANALYZE",  "Claude AI"),
    ]
    ko_labels = {
        "capture": "\uad11\uace0 \uc18c\uc7ac \uc2a4\ud06c\ub9b0\uc0f7 \uce90\uccad \uc911",
        "analyze": "Claude AI \ub9c8\ucf00\ud305 \ubd84\uc11d \uc911",
        "done":    "\ubd84\uc11d \uc644\ub8cc",
    }
    current_label = ko_labels.get(step, "\ucc98\ub9ac \uc911")

    steps_html = ""
    for i, (num, lbl, sub) in enumerate(step_defs):
        done_step   = i < step_idx
        active_step = (i == 0 and step == "capture") or (i == 2 and step == "analyze")

        if done_step:
            bg, bd, col, icon, fs, lbl_col = OK_DIM, f"1.5px solid {OK_BD}", OK, "v", "16px", OK
        elif active_step:
            bg, bd, col, icon, fs, lbl_col = AC_MED, f"2px solid {AC}", AC, "o", "13px", AC
        else:
            bg = "rgba(255,255,255,.03)"
            bd, col, icon, fs, lbl_col = f"1px solid {BORDER}", MU, num, "11px", MU

        ring = "animation:aia-ring 1.5s ease-in-out infinite;" if active_step else ""

        steps_html += (
            f'<div style="display:flex;flex-direction:column;align-items:center;gap:6px;flex:1;min-width:0;">'
            f'<div style="width:48px;height:48px;border-radius:50%;background:{bg};border:{bd};'
            f'display:flex;align-items:center;justify-content:center;'
            f'font-family:Pretendard,-apple-system,sans-serif;'
            f'font-size:{fs};font-weight:800;color:{col};flex-shrink:0;transition:.4s;{ring}">'
            f'{icon}</div>'
            f'<div style="font-size:9px;font-weight:700;color:{lbl_col};letter-spacing:1.5px;">{lbl}</div>'
            f'<div style="font-size:9px;color:{MU};letter-spacing:.5px;">{sub}</div>'
            f'</div>'
        )
        if i < len(step_defs) - 1:
            arr_col = OK if done_step else BORDER
            steps_html += (
                f'<div style="font-size:18px;color:{arr_col};padding-top:12px;flex-shrink:0;transition:.4s;">'
                f'&rsaquo;</div>'
            )

    dots_html = ""
    if not is_done:
        spans = "".join(
            f'<span style="width:4px;height:4px;border-radius:50%;background:{AC};opacity:.3;'
            f'display:inline-block;animation:aia-pulse 1.2s ease-in-out {i * 0.2:.1f}s infinite;">'
            f'</span>'
            for i in range(3)
        )
        dots_html = (
            f'<div style="display:flex;gap:6px;justify-content:center;margin-top:14px;">'
            f'{spans}</div>'
        )

    h_col  = OK if is_done else AC
    h_lbl  = "COMPLETE" if is_done else "ANALYZING"
    dot_a  = "" if is_done else "animation:aia-pulse 1.1s ease-in-out infinite;"
    bar_bg = OK if is_done else f"linear-gradient(90deg,#1a6dff,{AC})"
    bar_sh = f"box-shadow:0 0 8px {AC_GLOW};" if not is_done else f"box-shadow:0 0 6px rgba(61,255,160,.4);"
    bar_ani= "background-size:200% 100%;animation:aia-shimmer 1.8s linear infinite;" if not is_done else ""
    chk    = f'<span style="color:{OK};margin-right:6px;font-size:14px;">v</span>' if is_done else ""
    det    = (
        f'<div style="text-align:center;font-size:10px;color:{MU};'
        f'letter-spacing:.5px;margin-bottom:18px;">{detail}</div>'
    ) if detail else ""

    css = (
        "<style>"
        "@keyframes aia-scan{0%{transform:translateX(-100%);opacity:0}18%{opacity:1}82%{opacity:1}100%{transform:translateX(220%);opacity:0}}"
        "@keyframes aia-pulse{0%,100%{opacity:.3;transform:scale(1)}50%{opacity:1;transform:scale(1.5)}}"
        "@keyframes aia-ring{0%,100%{box-shadow:0 0 0 0 rgba(79,138,255,.45)}50%{box-shadow:0 0 0 8px rgba(79,138,255,0)}}"
        "@keyframes aia-shimmer{0%{background-position:0% 50%}100%{background-position:200% 50%}}"
        "@keyframes aia-rise{from{opacity:0;transform:translateY(10px)}to{opacity:1;transform:translateY(0)}}"
        "</style>"
    )

    return (
        css
        + f'<div style="background:{BG};border:1px solid {BORDER};border-radius:14px;'
        + f'padding:28px 32px;margin:4px 0 16px;position:relative;overflow:hidden;'
        + f'font-family:Pretendard,-apple-system,sans-serif;animation:aia-rise .35s cubic-bezier(.2,.8,.4,1);">'

        # Grid (same as splash)
        + f'<div style="position:absolute;inset:0;pointer-events:none;border-radius:14px;'
        + f'background-image:linear-gradient({GRID} 1px,transparent 1px),'
        + f'linear-gradient(90deg,{GRID} 1px,transparent 1px);'
        + f'background-size:48px 48px;opacity:.8;"></div>'

        # Scan line (same as splash)
        + f'<div style="position:absolute;top:0;left:0;right:0;height:2px;overflow:hidden;border-radius:14px 14px 0 0;">'
        + f'<div style="height:100%;background:linear-gradient(90deg,transparent,rgba(79,138,255,.55),rgba(91,157,255,.4),transparent);'
        + f'animation:aia-scan 2.3s ease-in-out infinite;"></div></div>'

        # Logo header (same as splash)
        + f'<div style="position:relative;display:flex;align-items:center;justify-content:space-between;margin-bottom:24px;">'
        + f'<div style="display:flex;align-items:center;gap:10px;">'
        + f'<span style="font-size:16px;font-weight:900;color:{TX};letter-spacing:-.5px;line-height:1;">'
        + f'AD<span style="color:{AC};">INTEL</span></span>'
        + f'<span style="font-size:9px;font-weight:600;color:{AC};border:1px solid {AC_BD};'
        + f'border-radius:4px;padding:2px 7px;letter-spacing:.5px;background:{AC_DIM};">v4.0</span>'
        + f'</div>'
        + f'<div style="display:flex;align-items:center;gap:8px;">'
        + f'<div style="width:7px;height:7px;border-radius:50%;background:{h_col};{dot_a}"></div>'
        + f'<span style="font-size:9px;font-weight:700;color:{h_col};letter-spacing:2.5px;">{h_lbl}</span>'
        + f'</div></div>'

        # Divider (same as splash)
        + f'<div style="position:relative;height:1px;background:linear-gradient(90deg,transparent,{AC_BD},transparent);margin-bottom:24px;"></div>'

        # Steps
        + f'<div style="position:relative;display:flex;align-items:flex-start;justify-content:center;gap:6px;margin-bottom:24px;">{steps_html}</div>'

        # Label
        + f'<div style="position:relative;text-align:center;margin-bottom:6px;">'
        + f'<span style="font-size:13px;font-weight:700;color:{TX};">{chk}{current_label}</span></div>'

        # Detail + dots
        + f'<div style="position:relative;">{det}{dots_html}</div>'

        # Divider
        + f'<div style="position:relative;height:1px;background:linear-gradient(90deg,transparent,{BORDER},transparent);margin:20px 0 16px;"></div>'

        # Progress bar (same as splash: 2px, glow, track)
        + f'<div style="position:relative;">'
        + f'<div style="background:{TRACK};border-radius:2px;height:2px;overflow:hidden;margin-bottom:10px;">'
        + f'<div style="height:100%;width:{pct}%;background:{bar_bg};border-radius:2px;{bar_sh}{bar_ani}transition:width .6s ease;"></div>'
        + f'</div>'
        + f'<div style="display:flex;justify-content:space-between;">'
        + f'<span style="font-size:9px;color:{MU};letter-spacing:1px;">{current} / {total} ASSETS</span>'
        + f'<span style="font-size:9px;color:{h_col};font-weight:700;">{pct}%</span>'
        + f'</div></div>'
        + f'</div>'
    )


def _safe_join(value, sep=", "):
    if value is None:
        return ""
    if isinstance(value, list):
        return sep.join(str(v) for v in value if v is not None and str(v).strip())
    if isinstance(value, dict):
        return sep.join(f"{k}: {v}" for k, v in value.items())
    return str(value)

def _normalize_ai_result(data: dict) -> dict:
    if not isinstance(data, dict):
        return {"_error": "AI 응답이 JSON object가 아닙니다."}
    vf = data.get("visual_facts") or {}
    ma = data.get("marketing_analysis") or {}
    scores = data.get("conversion_elements") or {}
    visible_text = vf.get("visible_text", data.get("copy", ""))
    objects = vf.get("objects", [])
    colors = vf.get("colors", [])
    appeal_type = ma.get("appeal_type", data.get("appeal", ""))
    data["copy"] = data.get("copy") or _safe_join(visible_text, " / ") or "확인 불가"
    data["hook"] = data.get("hook") or ma.get("hook_type") or "확인 불가"
    data["appeal"] = data.get("appeal") or _safe_join(appeal_type) or "확인 불가"
    data["tone"] = data.get("tone") or _safe_join(colors) or "확인 불가"
    data["target"] = data.get("target") or ma.get("target") or "확인 불가"
    data["message"] = data.get("message") or ma.get("message") or "확인 불가"
    data["evidence"] = data.get("evidence") or ma.get("evidence") or "확인 불가"
    data["main_visual"] = data.get("main_visual") or _safe_join(objects) or vf.get("product_focus") or "확인 불가"
    data["layout_type"] = data.get("layout_type") or vf.get("layout_type") or "확인 불가"
    data["scores"] = scores
    tags = data.get("tags") or []
    if not tags:
        base_tags = []
        for v in [appeal_type, vf.get("layout_type"), objects[:2] if isinstance(objects, list) else objects]:
            if isinstance(v, list):
                base_tags += [str(x) for x in v if x]
            elif v:
                base_tags.append(str(v))
        data["tags"] = base_tags[:5]
    return data

def _force_sanitize_text(text: str) -> str:
    if not isinstance(text, str):
        return text
    replacements = [
        ("곰모양", "브라운 컬러의 둥근 형태"), ("곰 모양", "브라운 컬러의 둥근 형태"),
        ("곰 캐릭터", "브라운 컬러의 둥근 형태 비주얼"), ("곰", "브라운 컬러의 둥근 형태"),
        ("동물 캐릭터", "형태 비주얼"), ("캐릭터 동물", "형태 비주얼"),
        ("동물", "형태 비주얼"), ("토끼", "형태 비주얼"),
        ("강아지", "형태 비주얼"), ("고양이", "형태 비주얼"),
    ]
    out = text
    for bad, repl in replacements:
        out = out.replace(bad, repl)
    return out

def _sanitize_json_deep(value):
    if isinstance(value, dict):
        return {k: _sanitize_json_deep(v) for k, v in value.items()}
    if isinstance(value, list):
        return [_sanitize_json_deep(v) for v in value]
    if isinstance(value, str):
        return _force_sanitize_text(value)
    return value

def _extract_ocr_with_claude(img_b64: str) -> dict:
    if not API_KEY or not img_b64:
        return {"texts": [], "raw": "", "_error": "OCR 입력 이미지 없음"}
    try:
        ocr_prompt = """
당신은 광고 이미지 OCR 추출기입니다.
이미지 안에 실제로 보이는 글자만 최대한 그대로 추출하세요.
규칙:
- 보이는 텍스트만 추출, 해석 금지
- 글자가 불확실하면 "확인 불가"
- JSON 외 설명문/마크다운/백틱 금지
- 반드시 json.loads() 가능한 순수 JSON만 반환

반환 형식:
{"texts":["이미지 안 실제 문구1","이미지 안 실제 문구2"],"brand_candidates":["보이는 브랜드명"],"product_candidates":["보이는 제품명"],"price_texts":["가격/할인 관련 문구"],"cta_texts":["CTA 문구"]}
"""
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"x-api-key": API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            json={
                "model": _anthropic_model(), "max_tokens": 700, "temperature": 0,
                "messages": [{"role": "user", "content": [
                    {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": img_b64}},
                    {"type": "text", "text": ocr_prompt},
                ]}],
            },
            timeout=45,
        )
        if resp.status_code != 200:
            return {"texts": [], "raw": "", "_error": f"OCR API {resp.status_code}: {resp.text[:200]}"}
        txt = resp.json()["content"][0]["text"].strip()
        parsed, err = _parse_ai_json(txt)
        if err:
            return {"texts": [], "raw": txt, "_error": err}
        if not isinstance(parsed, dict):
            return {"texts": [], "raw": txt, "_error": "OCR 응답이 JSON object가 아닙니다."}
        for k in ["texts", "brand_candidates", "product_candidates", "price_texts", "cta_texts"]:
            v = parsed.get(k, [])
            if isinstance(v, str):
                parsed[k] = [v]
            elif not isinstance(v, list):
                parsed[k] = []
        parsed["raw"] = txt
        return _sanitize_json_deep(parsed)
    except Exception as ex:
        return {"texts": [], "raw": "", "_error": f"OCR 실패: {str(ex)[:200]}"}

def analyze(item):
    """
    근거 기반 광고 소재 분석.
    철학: 보이는 것(OBSERVE) → 인용해서 해석(INTERPRET).
    보이지 않으면 확인 불가. 추정·추론·가정 금지.
    """
    if not API_KEY:
        return {"_error": "API Key 없음"}
    b64 = item.get("img_b64", "") or db_get_img_b64(item["id"])
    if not b64:
        return {"_error": "스크린샷 확보 실패 — 선택 분석 버튼으로 재시도하세요."}

    capture_src = item.get("capture_source", "unknown")
    has_card = "card" in capture_src

    try:
        # ── 1차: OCR 추출 ──
        ocr_data     = _extract_ocr_with_claude(b64)
        ocr_texts    = ocr_data.get("texts", [])
        brand_text   = _safe_join(ocr_data.get("brand_candidates", []), " / ")
        product_text = _safe_join(ocr_data.get("product_candidates", []), " / ")
        price_text   = _safe_join(ocr_data.get("price_texts", []), " / ")
        cta_text     = _safe_join(ocr_data.get("cta_texts", []), " / ")
        ocr_all      = _safe_join(ocr_texts, " / ")

        card_note = (
            "캡처 범위: 광고 카드 전체 (이미지 + 텍스트 카피 + CTA + 브랜드 영역 포함 가능)"
            if has_card else
            "캡처 범위: 광고 이미지 영역 단독"
        )

        prompt = f"""당신은 Meta 광고 소재 분석 전문가입니다.
아래 두 단계를 순서대로 수행하세요.

━━━ 분석 컨텍스트 ━━━
검색 키워드: {item["keyword"]}
{card_note}

━━━ OCR 사전 추출 결과 ━━━
(아래는 별도 OCR 모델이 이미지에서 추출한 텍스트입니다. 이미지와 대조하여 검증하세요.)
전체 문구: {ocr_all or "(추출된 텍스트 없음)"}
브랜드 후보: {brand_text or "(없음)"}
제품명 후보: {product_text or "(없음)"}
가격/할인: {price_text or "(없음)"}
CTA: {cta_text or "(없음)"}

━━━ STEP 1: OBSERVE — 관찰만 하세요 ━━━
이미지에서 실제로 보이는 것만 기록합니다.
규칙:
- 텍스트: 이미지에서 눈으로 읽히는 문자를 그대로 옮기세요. 흐리거나 잘려서 읽을 수 없으면 기록하지 마세요.
- 비주얼: "무엇처럼 보인다" 금지. "색상 X, 위치 Y, 크기 Z" 형태로 객관적으로 기술하세요.
- 추측·추정·가정 완전 금지. 보이지 않으면 해당 필드를 빈 배열 또는 false로 두세요.
- "확인 불가"는 정직한 답변입니다. 보이지 않는 것을 채우지 마세요.

━━━ STEP 2: INTERPRET — STEP 1 근거로만 해석하세요 ━━━
규칙:
- 모든 마케팅 해석은 반드시 STEP 1에서 확인된 요소를 evidence로 인용해야 합니다.
- STEP 1에 없는 정보로 해석하면 안 됩니다.
- 타겟 추정: 이미지에서 인물, 제품 카테고리, 텍스트 단서가 있을 때만 가능합니다.
  단서가 없으면 "확인 불가"가 올바른 답입니다.
- 후크/소구 판단: 이미지에서 관찰된 레이아웃과 텍스트를 근거로만 판단하세요.
- 전환력 점수: 실제로 보이는 요소 기준으로 채점하세요. 없으면 1점, 강하면 5점.

━━━ JSON 출력 필수 규칙 ━━━
- 순수 JSON만 반환. 마크다운·백틱·설명문 절대 금지.
- 문자열 값 내부에 실제 줄바꿈(개행) 절대 금지 — 모든 텍스트는 한 줄로.
- 문자열 값 내부의 큰따옴표(")는 반드시 작은따옴표(')로 교체.
- 배열 요소 사이에 실제 줄바꿈 금지 — 쉼표로만 구분.
- 마지막 요소 뒤 쉼표(trailing comma) 금지.
- json.loads()로 파싱 가능해야 합니다.

━━━ JSON 구조 ━━━
{{
  "observe": {{
    "text_exact": ["이미지에서 실제로 읽힌 문구를 그대로. 못 읽으면 빈 배열"],
    "layout_description": "이미지 구성을 좌/우/상/하 위치 기준으로 객관적으로 기술",
    "main_visual_description": "주요 피사체를 색상·형태·위치만으로 기술. 정체 추정 금지",
    "supporting_visuals": ["부차적 요소들 — 색상·형태·위치만"],
    "color_palette": ["배경색", "주요 요소 색상"],
    "human_present": false,
    "human_description": "사람이 있으면 외형만 기술. 없으면 null",
    "price_text_exact": "가격·할인 텍스트 그대로. 없으면 null",
    "cta_text_exact": "CTA 버튼·문구 그대로. 없으면 null",
    "brand_text_exact": "브랜드명 텍스트 그대로. 없으면 null",
    "layout_type": "제품단독형/모델사용형/Before-After형/리뷰UGC형/가격혜택형/카드뉴스형/문제제기형/기능설명형/브랜드필름형"
  }},
  "interpret": {{
    "hook_type": "observe에서 확인된 요소 기반. 근거 없으면 확인 불가",
    "hook_evidence": "후크 판단의 근거 — observe 필드를 인용",
    "primary_appeal": "observe에서 확인된 요소 기반. 근거 없으면 확인 불가",
    "appeal_evidence": "소구 판단의 근거 — observe 필드를 인용",
    "target_audience": "observe의 인물·제품·텍스트 단서 기반. 단서 없으면 확인 불가",
    "target_evidence": "타겟 판단의 근거. 단서 없으면 null",
    "core_message": "observe의 텍스트·비주얼 기반 메시지. 텍스트가 없으면 확인 불가",
    "emotional_trigger": "observe 근거 있을 때만. 없으면 확인 불가"
  }},
  "scores": {{
    "price_emphasis": 1,
    "product_visibility": 1,
    "readability": 1,
    "discount_visibility": 1,
    "visual_clarity": 1,
    "cta_strength": 1,
    "overall_conversion_power": 1,
    "score_rationale": "observe에서 확인된 요소 기준으로 채점 근거 기술"
  }},
  "diagnosis": {{
    "strengths": ["observe 근거 있는 강점만"],
    "weaknesses": ["observe 근거 있는 약점만"],
    "improvement_direction": "관찰된 약점 기반 개선 방향"
  }},
  "tags": ["관찰 근거 있는 태그만", "최대 5개"]
}}

점수: 1=없음/매우약함 2=약함 3=보통 4=강함 5=매우강함
보이지 않는 요소는 1점입니다.
"""

        payload = [
            {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": b64}},
            {"type": "text", "text": prompt},
        ]
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"x-api-key": API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            json={
                "model": _anthropic_model("analysis"),
                "max_tokens": 1800,
                "temperature": 0,
                "messages": [{"role": "user", "content": payload}],
            },
            timeout=60,
        )
        if resp.status_code != 200:
            return {"_error": f"API {resp.status_code}: {resp.text[:300]}"}

        txt = resp.json()["content"][0]["text"].strip()
        parsed, parse_error = _parse_ai_json(txt)
        if parse_error:
            return {"_error": parse_error}

        normalized = _normalize_observe_interpret(parsed)
        normalized = _sanitize_json_deep(normalized)
        normalized["ocr"] = _sanitize_json_deep(ocr_data)
        return normalized

    except Exception as ex:
        return {"_error": str(ex)[:300]}


def _normalize_observe_interpret(data: dict) -> dict:
    """
    새 observe/interpret 구조를 UI 호환 top-level 필드로 매핑.
    추정·fallback 없음 — 없으면 그냥 없는 것.
    """
    if not isinstance(data, dict):
        return {"_error": "AI 응답이 JSON object가 아닙니다."}

    ob = data.get("observe") or {}
    ip = data.get("interpret") or {}
    sc = data.get("scores") or {}
    dg = data.get("diagnosis") or {}

    data["copy"]        = _safe_join(ob.get("text_exact", []), " / ") or "확인 불가"
    data["hook"]        = ip.get("hook_type") or "확인 불가"
    data["appeal"]      = ip.get("primary_appeal") or "확인 불가"
    data["tone"]        = _safe_join(ob.get("color_palette", [])) or "확인 불가"
    data["target"]      = ip.get("target_audience") or "확인 불가"
    data["message"]     = ip.get("core_message") or "확인 불가"
    data["evidence"]    = ip.get("hook_evidence") or ip.get("appeal_evidence") or "확인 불가"
    data["main_visual"] = ob.get("main_visual_description") or "확인 불가"
    data["layout_type"] = ob.get("layout_type") or "확인 불가"
    data["scores"]      = sc

    # visual_facts 호환
    data["visual_facts"] = {
        "visible_text":  ob.get("text_exact", []),
        "layout_type":   ob.get("layout_type", "확인 불가"),
        "ocr_text":      ob.get("text_exact", []),
    }

    # marketing_analysis 호환
    data["marketing_analysis"] = {
        "hook_type":    ip.get("hook_type", "확인 불가"),
        "appeal_type":  [ip.get("primary_appeal", "확인 불가")],
        "target":       ip.get("target_audience", "확인 불가"),
        "message":      ip.get("core_message", "확인 불가"),
        "evidence":     ip.get("hook_evidence", "확인 불가"),
    }

    # conversion_elements 호환
    data["conversion_elements"] = {
        **sc,
        "score_reason": sc.get("score_rationale", ""),
    }

    # creative_diagnosis 호환
    data["creative_diagnosis"] = {
        "strengths":             dg.get("strengths", []),
        "weaknesses":            dg.get("weaknesses", []),
        "improvement_direction": dg.get("improvement_direction", "확인 불가"),
    }

    if not data.get("tags"):
        data["tags"] = []

    return data


def analyze_parallel(items, max_workers=6, force=False):
    if not API_KEY:
        return items, []
    id_map = {item["id"]: item for item in items}
    errors = []

    def task(item):
        return item["id"], analyze(item)

    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        futures = {
            ex.submit(task, item): item["id"]
            for item in items
            if force or not item.get("ai")
        }
        for future in as_completed(futures):
            aid = futures[future]
            try:
                aid, result = future.result(timeout=90)
                if result and aid in id_map:
                    id_map[aid]["ai"] = result
                    db_set_field(aid, ai_json=json.dumps(result, ensure_ascii=False))
                    if result.get("_error"):
                        errors.append(f"[{id_map[aid]['keyword']}] {result['_error']}")
            except Exception as ex_inner:
                err_msg = f"[{aid[:8]}] 분석 예외: {str(ex_inner)[:120]}"
                errors.append(err_msg)

    return items, errors

def test_api():
    if not API_KEY:
        return False, "API Key가 없습니다. Streamlit Secrets를 확인하세요."
    try:
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"x-api-key": API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            json={"model": "claude-haiku-4-5-20251001", "max_tokens": 10, "messages": [{"role": "user", "content": "hi"}]},
            timeout=10,
        )
        if resp.status_code == 200:
            return True, "API 정상 연결됨"
        return False, f"API 오류 {resp.status_code}: {resp.text[:300]}"
    except Exception as ex:
        return False, f"연결 실패: {ex}"

def summarize_insights(analyzed_items):
    if not API_KEY or not analyzed_items:
        return None
    snippets = []
    for a in analyzed_items[:20]:
        ai = a.get("ai") or {}
        if not ai:
            continue
        snippets.append(
            "- 소구:" + ai.get("appeal", "") +
            " / 후크:" + ai.get("hook", "") +
            " / 타겟:" + ai.get("target", "") +
            " / 메시지:" + ai.get("message", "") +
            " / 약점:" + _safe_join((ai.get("creative_diagnosis") or {}).get("weaknesses", []))
        )
    if not snippets:
        return None
    json_schema = (
        '{"dominant_appeal":"가장 많이 쓰인 소구 유형 + 비율 설명",'
        '"common_target":"공통 타겟 고객 요약",'
        '"key_message":"반복되는 핵심 메시지 패턴",'
        '"strategy":"이 광고들이 공유하는 전략적 방향 (2~3문장)",'
        '"common_weakness":"경쟁사들이 공통으로 놓친 약점 (2문장)",'
        '"our_opportunity":"경쟁 우위를 위한 차별화 기회 (2~3문장)",'
        '"recommendations":"구체적인 소재 제작 제안 (2~3문장)",'
        '"tags":["태그1","태그2","태그3","태그4","태그5"]}'
    )
    prompt = (
        "광고 전략 전문가입니다. 아래는 Meta 광고 소재 분석 결과들입니다.\n\n" +
        "\n".join(snippets) +
        "\n\n이 소재들을 종합해 아래 JSON만 반환 (마크다운/백틱 없이 순수 JSON):\n" + json_schema
    )
    try:
        resp = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"x-api-key": API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            json={"model": "claude-haiku-4-5-20251001", "max_tokens": 900,
                  "messages": [{"role": "user", "content": prompt}]},
            timeout=30,
        )
        if resp.status_code != 200:
            return None
        txt = resp.json()["content"][0]["text"].strip()
        parsed, _ = _parse_ai_json(txt)
        return parsed
    except Exception:
        return None


# ============================================================
# 유틸
# ============================================================
def merge(existing, new_items):
    known = {make_fp(a) for a in existing}
    known |= {make_fp(a) for a in existing if a["id"] in st.session_state.hidden}
    out, n = list(existing), 0
    for item in new_items:
        k = make_fp(item)
        if k in known:
            continue
        known.add(k)
        out.append(item)
        n += 1
    return out, n

def get_visible(ftype, fkw, fstar, fai):
    hidden = st.session_state.hidden
    items = [a for a in st.session_state.assets if a["id"] not in hidden]
    if ftype != "전체":
        items = [a for a in items if a["asset_type"] == ftype]
    if fkw and fkw != "전체":
        items = [a for a in items if a["keyword"] == fkw]
    if fstar:
        items = [a for a in items if a.get("starred")]
    if fai:
        items = [a for a in items if a.get("ai")]
    return items

def do_reset():
    db_reset()
    st.session_state.assets   = []
    st.session_state.hidden   = set()
    st.session_state.history  = []
    st.session_state.log      = []
    st.session_state.selected = set()
    st.session_state.summary  = None

def toggle_star(aid):
    for a in st.session_state.assets:
        if a["id"] == aid:
            a["starred"] = not a.get("starred", False)
            db_set_field(aid, starred=1 if a["starred"] else 0)
            break

def to_csv(items):
    import csv
    fields = ["keyword", "country", "asset_type", "image_url", "source_url",
              "caption", "width", "height", "created_at", "starred"]
    buf = io.StringIO()
    w = csv.DictWriter(buf, fieldnames=fields)
    w.writeheader()
    for a in items:
        w.writerow({f: a.get(f, "") for f in fields})
    return buf.getvalue().encode("utf-8-sig")

def fetch_image_bytes(url: str) -> bytes | None:
    try:
        resp = requests.get(url, timeout=8, headers={
            "User-Agent": (
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0.0.0 Safari/537.36"
            ),
            "Referer": "https://www.facebook.com/",
        })
        if resp.status_code == 200 and len(resp.content) > 1000:
            return resp.content
    except Exception:
        pass
    try:
        ensure_browser()
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True, args=["--no-sandbox"])
            page = browser.new_page(viewport={"width": 800, "height": 800})
            page.goto(url, wait_until="domcontentloaded", timeout=15000)
            page.wait_for_timeout(1500)
            img_bytes = page.screenshot(type="jpeg", quality=80)
            browser.close()
            return img_bytes
    except Exception:
        return None


# ============================================================
# PPT 내보내기 — 고도화 (경쟁 현황 + 기회 슬라이드 추가)
# ============================================================
def to_pptx(items: list, summary: dict | None = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import tempfile
    from collections import Counter

    # 색상 팔레트
    C_BG    = RGBColor(0xFF, 0xFF, 0xFF)
    C_DARK  = RGBColor(0x11, 0x13, 0x18)
    C_BLUE  = RGBColor(0x1A, 0x6D, 0xFF)
    C_GRAY  = RGBColor(0x3A, 0x3F, 0x4A)
    C_MUTED = RGBColor(0x8A, 0x90, 0x9E)
    C_BG2   = RGBColor(0xF7, 0xF8, 0xFA)
    C_GREEN = RGBColor(0x00, 0xA8, 0x6B)
    C_ORANGE= RGBColor(0xE0, 0x7B, 0x00)
    C_BDBLUE= RGBColor(0xCC, 0xDD, 0xFF)
    C_NAVY  = RGBColor(0x0A, 0x1A, 0x40)

    def add_text(slide, text, x, y, w, h, size=12, bold=False,
                 color=None, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or C_DARK
        return txBox

    def add_rect(slide, x, y, w, h, fill_color, line_color=None, transparency=0):
        shape = slide.shapes.add_shape(1,
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if transparency:
            shape.fill.fore_color.theme_color = None
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(0.5)
        else:
            shape.line.fill.background()
        return shape

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── 1. 표지 슬라이드 ──
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, C_NAVY)
    add_rect(slide, 0, 0, 4.0, 7.5, C_DARK)
    add_rect(slide, 4.0, 0, 0.06, 7.5, C_BLUE)

    add_text(slide, "AD INTEL",               4.3, 1.6, 8, 1.2, size=44, bold=True, color=C_BG)
    add_text(slide, "META AD LIBRARY INTELLIGENCE REPORT",
                                               4.3, 2.9, 8, 0.6, size=12, color=C_BLUE)
    kw_list = list(dict.fromkeys(a["keyword"] for a in items))
    add_text(slide, "분석 키워드: " + " · ".join(kw_list),
                                               4.3, 3.6, 8, 0.5, size=11, color=C_MUTED)
    add_text(slide, f"소재 {len(items)}개  ·  생성일 {time.strftime('%Y.%m.%d')}",
                                               4.3, 4.1, 8, 0.5, size=11, color=C_MUTED)

    analyzed_count = sum(1 for a in items if a.get("ai") and not (a.get("ai") or {}).get("_error"))
    add_text(slide, f"AI 분석 완료: {analyzed_count}개",
                                               4.3, 4.6, 8, 0.4, size=10, color=C_GREEN)

    # 목차
    toc_items = ["경쟁 현황 분석", "소구 유형 분포", "소재별 상세 분석"]
    if summary:
        toc_items += ["총합 인사이트", "기회 영역 & 제안"]
    for i, t in enumerate(toc_items):
        add_rect(slide, 0.3, 1.5 + i * 0.8, 3.2, 0.6, C_BLUE if i == 0 else RGBColor(0x1A, 0x22, 0x50))
        add_text(slide, f"{i+1:02d}  {t}", 0.5, 1.6 + i * 0.8, 3.0, 0.4,
                 size=10, bold=(i == 0), color=C_BG)

    # ── 2. 경쟁 현황 요약 슬라이드 ──
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG)
    add_rect(slide, 0, 0, 13.33, 1.1, C_DARK)

    add_text(slide, "01  경쟁 현황 분석", 0.4, 0.2, 12, 0.6, size=18, bold=True, color=C_BG)
    add_text(slide, f"총 {len(items)}개 소재 · AI 분석 {analyzed_count}개 기준",
             0.4, 0.72, 12, 0.32, size=10, color=C_MUTED)

    # 통계 카드 4개
    ai_items = [a for a in items if a.get("ai") and not (a.get("ai") or {}).get("_error")]
    hooks = [a["ai"].get("hook", "확인 불가") for a in ai_items]
    appeals = []
    for a in ai_items:
        ap = a["ai"].get("appeal", "")
        if isinstance(ap, list):
            appeals += ap
        elif ap:
            appeals.append(ap)
    layouts = [((a["ai"].get("visual_facts") or {}).get("layout_type") or a["ai"].get("layout_type") or "확인 불가")
               for a in ai_items]
    scores_all = [((a["ai"].get("conversion_elements") or {}).get("overall_conversion_power") or 0)
                  for a in ai_items]
    avg_score = round(sum(scores_all) / len(scores_all), 1) if scores_all else 0

    top_hook   = Counter(hooks).most_common(1)[0][0] if hooks else "—"
    top_appeal = Counter(appeals).most_common(1)[0][0] if appeals else "—"
    top_layout = Counter(layouts).most_common(1)[0][0] if layouts else "—"

    stat_cards = [
        ("TOTAL ASSETS", str(len(items)), "수집된 소재 수"),
        ("TOP HOOK", top_hook, "가장 많이 쓰인 후크"),
        ("TOP APPEAL", top_appeal, "주요 소구 유형"),
        ("AVG CONVERSION", f"{avg_score} / 5", "평균 전환력 점수"),
    ]
    for i, (lbl, val, sub) in enumerate(stat_cards):
        cx = 0.3 + i * 3.25
        add_rect(slide, cx, 1.3, 3.0, 1.6, C_BG2)
        add_rect(slide, cx, 1.3, 0.06, 1.6, C_BLUE)
        add_text(slide, lbl, cx + 0.2, 1.4, 2.7, 0.35, size=8, bold=True, color=C_BLUE)
        add_text(slide, val, cx + 0.2, 1.8, 2.7, 0.7, size=20, bold=True, color=C_DARK)
        add_text(slide, sub, cx + 0.2, 2.6, 2.7, 0.25, size=9, color=C_MUTED)

    # 소재 타입 구성
    img_cnt = sum(1 for a in items if a["asset_type"] == "image")
    vid_cnt = sum(1 for a in items if a["asset_type"] == "video_poster")
    add_text(slide, "소재 타입 구성", 0.3, 3.2, 6, 0.35, size=11, bold=True, color=C_DARK)
    add_rect(slide, 0.3, 3.6, (img_cnt / max(len(items), 1)) * 6.0, 0.4, C_BLUE)
    add_rect(slide, 0.3 + (img_cnt / max(len(items), 1)) * 6.0, 3.6,
             (vid_cnt / max(len(items), 1)) * 6.0, 0.4, C_ORANGE)
    add_text(slide, f"이미지 {img_cnt}개 ({img_cnt*100//max(len(items),1)}%)",
             0.3, 4.05, 3, 0.3, size=9, color=C_BLUE)
    add_text(slide, f"비디오 {vid_cnt}개 ({vid_cnt*100//max(len(items),1)}%)",
             3.5, 4.05, 3, 0.3, size=9, color=C_ORANGE)

    # 키워드별 소재 수
    kw_counts = Counter(a["keyword"] for a in items)
    add_text(slide, "키워드별 소재 수", 7.0, 3.2, 6, 0.35, size=11, bold=True, color=C_DARK)
    y_kw = 3.6
    for kw, cnt in kw_counts.most_common(5):
        bar_w = (cnt / max(len(items), 1)) * 5.5
        add_rect(slide, 7.0, y_kw, bar_w, 0.3, C_BLUE)
        add_text(slide, f"{kw}  ({cnt})", 7.0 + bar_w + 0.1, y_kw, 2.5, 0.3,
                 size=9, color=C_DARK)
        y_kw += 0.42

    # 분석 에러 현황
    err_count = sum(1 for a in items if (a.get("ai") or {}).get("_error"))
    if err_count:
        add_rect(slide, 0.3, 6.8, 12.7, 0.4, RGBColor(0xFF, 0xF0, 0xF0))
        add_text(slide, f"⚠  분석 오류 {err_count}개 포함 (캡처 실패 또는 API 오류)",
                 0.5, 6.85, 12, 0.3, size=9, color=C_ORANGE)

    # ── 3. 소구/후크 분포 슬라이드 ──
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG)
    add_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
    add_text(slide, "02  소구 & 후크 유형 분포", 0.4, 0.2, 12, 0.6, size=18, bold=True, color=C_BG)
    add_text(slide, f"AI 분석 {analyzed_count}개 기준", 0.4, 0.72, 12, 0.32, size=10, color=C_MUTED)

    # 소구 분포 (좌)
    add_text(slide, "소구 유형 분포", 0.4, 1.25, 6, 0.35, size=12, bold=True, color=C_DARK)
    appeal_cnt = Counter(appeals).most_common(7)
    max_ap = appeal_cnt[0][1] if appeal_cnt else 1
    colors_ap = [C_BLUE, RGBColor(0x4F, 0x8A, 0xFF), RGBColor(0x7C, 0x5C, 0xFF),
                 RGBColor(0x00, 0xA8, 0x6B), C_ORANGE,
                 RGBColor(0xFF, 0x6B, 0x6B), RGBColor(0xFF, 0xB8, 0x4F)]
    for i, (ap, cnt) in enumerate(appeal_cnt):
        y_pos = 1.7 + i * 0.56
        bar_w = (cnt / max(max_ap, 1)) * 5.0
        add_rect(slide, 1.8, y_pos + 0.05, bar_w, 0.35,
                 colors_ap[i % len(colors_ap)])
        add_text(slide, ap, 0.4, y_pos, 1.3, 0.4, size=9, color=C_DARK)
        add_text(slide, str(cnt), 1.8 + bar_w + 0.1, y_pos, 0.5, 0.4, size=9,
                 bold=True, color=C_DARK)

    # 후크 분포 (우)
    add_text(slide, "후크 유형 분포", 7.0, 1.25, 6, 0.35, size=12, bold=True, color=C_DARK)
    hook_cnt = Counter(hooks).most_common(7)
    max_hk = hook_cnt[0][1] if hook_cnt else 1
    for i, (hk, cnt) in enumerate(hook_cnt):
        y_pos = 1.7 + i * 0.56
        bar_w = (cnt / max(max_hk, 1)) * 5.0
        add_rect(slide, 8.7, y_pos + 0.05, bar_w, 0.35, colors_ap[i % len(colors_ap)])
        add_text(slide, hk, 7.0, y_pos, 1.6, 0.4, size=9, color=C_DARK)
        add_text(slide, str(cnt), 8.7 + bar_w + 0.1, y_pos, 0.5, 0.4, size=9,
                 bold=True, color=C_DARK)

    # 레이아웃 분포 (하단)
    add_text(slide, "레이아웃 유형 분포", 0.4, 6.1, 12, 0.35, size=11, bold=True, color=C_DARK)
    layout_cnt = Counter(layouts).most_common(6)
    total_l = sum(c for _, c in layout_cnt)
    x_l = 0.4
    for lyt, cnt in layout_cnt:
        w_l = (cnt / max(total_l, 1)) * 12.0
        if w_l < 0.3:
            continue
        add_rect(slide, x_l, 6.5, w_l - 0.05, 0.5, C_BLUE)
        if w_l > 1.2:
            add_text(slide, f"{lyt[:6]}\n{cnt}개", x_l + 0.05, 6.5, w_l - 0.1, 0.5,
                     size=7, color=C_BG, wrap=True)
        x_l += w_l

    # ── 4. 소재별 슬라이드 ──
    for idx, item in enumerate(items):
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, C_BG)
        add_rect(slide, 0, 0, 13.33, 0.55, C_DARK)
        add_text(slide, f"03  소재 상세  ·  {idx+1} / {len(items)}",
                 0.3, 0.08, 9, 0.38, size=11, bold=True, color=C_BG)
        add_text(slide, item["keyword"].upper() + "  ·  " + item["country"],
                 10.0, 0.08, 3, 0.38, size=10, color=C_BLUE, align=PP_ALIGN.RIGHT)

        # 이미지
        img_bytes = fetch_image_bytes(item["image_url"])
        if img_bytes:
            with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
                tmp.write(img_bytes)
                tmp_path = tmp.name
            try:
                pic = slide.shapes.add_picture(tmp_path, Inches(0.3), Inches(0.7), width=Inches(4.2))
                if pic.height > Inches(6.2):
                    pic.height = Inches(6.2)
                    pic.width  = Inches(4.2)
            except Exception:
                add_rect(slide, 0.3, 0.7, 4.2, 6.2, C_BG2)
                add_text(slide, "이미지 로드 실패", 0.3, 3.5, 4.2, 0.5,
                         align=PP_ALIGN.CENTER, color=C_MUTED)
            import os
            os.unlink(tmp_path)
        else:
            add_rect(slide, 0.3, 0.7, 4.2, 6.2, C_BG2)
            add_text(slide, "이미지 없음", 0.3, 3.5, 4.2, 0.5,
                     align=PP_ALIGN.CENTER, color=C_MUTED)

        # 세로 구분선
        add_rect(slide, 4.7, 0.55, 0.02, 6.9, C_BDBLUE)

        rx, rw = 5.0, 8.1

        ai = item.get("ai")
        if ai and not ai.get("_error"):
            vf = ai.get("visual_facts") or {}
            ce = ai.get("conversion_elements") or {}
            cd = ai.get("creative_diagnosis") or {}
            ma = ai.get("marketing_analysis") or {}

            # 분석 필드
            fields = [
                ("후크",    ai.get("hook") or ma.get("hook_type") or "—"),
                ("소구",    ai.get("appeal") or "—"),
                ("타겟",    ai.get("target") or "—"),
                ("레이아웃", vf.get("layout_type") or ai.get("layout_type") or "—"),
                ("메시지",  ai.get("message") or "—"),
                ("근거",    ai.get("evidence") or "—"),
            ]
            y_f = 0.65
            for lbl, val in fields:
                add_rect(slide, rx, y_f, 1.0, 0.32, C_BG2)
                add_text(slide, lbl, rx + 0.05, y_f, 0.95, 0.32, size=8, bold=True, color=C_BLUE)
                val_str = str(val)[:100]
                add_text(slide, val_str, rx + 1.1, y_f, rw - 1.15, 0.38, size=9, color=C_DARK, wrap=True)
                y_f += 0.48

            # 전환력 점수 바
            add_text(slide, "전환력 점수", rx, y_f + 0.1, rw, 0.3, size=9, bold=True, color=C_DARK)
            score_items = [
                ("가격강조", ce.get("price_emphasis", 0)),
                ("제품가시성", ce.get("product_visibility", 0)),
                ("가독성", ce.get("readability", 0)),
                ("할인가시성", ce.get("discount_visibility", 0)),
                ("시각명확도", ce.get("visual_clarity", 0)),
                ("전환력", ce.get("overall_conversion_power", 0)),
            ]
            y_sc = y_f + 0.45
            for s_lbl, s_val in score_items:
                s_val = int(s_val) if s_val else 0
                add_text(slide, s_lbl, rx, y_sc, 1.3, 0.28, size=7, color=C_MUTED)
                bar_total = 1.3
                for b in range(5):
                    add_rect(slide, rx + 1.35 + b * (bar_total / 5 + 0.02),
                             y_sc + 0.04, bar_total / 5, 0.2,
                             C_BLUE if b < s_val else C_BG2)
                add_text(slide, str(s_val), rx + 1.35 + 5 * (bar_total / 5 + 0.02) + 0.05,
                         y_sc, 0.3, 0.28, size=7, bold=True, color=C_DARK)
                y_sc += 0.3

            # 장점 / 약점 / 개선
            y_dia = y_sc + 0.1
            if y_dia < 6.5:
                strengths = _safe_join(cd.get("strengths", []))
                weaknesses = _safe_join(cd.get("weaknesses", []))
                improve = cd.get("improvement_direction", "—")
                for emoji, lbl, val in [("✓", "장점", strengths), ("✕", "약점", weaknesses), ("→", "개선", improve)]:
                    add_text(slide, f"{emoji} {lbl}", rx, y_dia, 1.0, 0.3, size=8, bold=True, color=C_GRAY)
                    add_text(slide, str(val)[:120], rx + 1.0, y_dia, rw - 1.05, 0.35, size=8, color=C_DARK, wrap=True)
                    y_dia += 0.38

            # 태그
            tags = ai.get("tags", [])
            if tags:
                tag_str = "  ".join(f"#{t}" for t in tags[:6])
                add_text(slide, tag_str, rx, 7.1, rw, 0.3, size=9, color=C_BLUE)
        elif ai and ai.get("_error"):
            add_rect(slide, rx, 1.0, rw, 0.6, RGBColor(0xFF, 0xF0, 0xF0))
            add_text(slide, "⚠  분석 오류: " + ai["_error"][:120],
                     rx + 0.1, 1.1, rw - 0.2, 0.4, size=9, color=C_ORANGE, wrap=True)
        else:
            add_text(slide, "AI 분석 없음\n(선택 분석 실행 후 PPT 재생성 권장)",
                     rx, 2.5, rw, 1.0, size=11, color=C_MUTED)

    # ── 5. 총합 인사이트 슬라이드 ──
    if summary:
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, C_NAVY)
        add_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
        add_text(slide, "04  총합 인사이트", 0.4, 0.2, 12, 0.6, size=18, bold=True, color=C_BG)
        add_text(slide, "선택 소재 전체 종합 분석", 0.4, 0.72, 12, 0.32, size=10, color=C_MUTED)

        grid_items = [
            ("DOMINANT APPEAL", summary.get("dominant_appeal", "—")),
            ("COMMON TARGET",   summary.get("common_target",   "—")),
            ("KEY MESSAGE",     summary.get("key_message",     "—")),
            ("COMMON WEAKNESS", summary.get("common_weakness", "—")),
        ]
        for i, (lbl, val) in enumerate(grid_items):
            col = i % 2
            row = i // 2
            gx = 0.4 + col * 6.4
            gy = 1.3 + row * 1.55
            add_rect(slide, gx, gy, 6.0, 1.4, RGBColor(0x1A, 0x20, 0x30))
            add_rect(slide, gx, gy, 0.06, 1.4, C_BLUE)
            add_text(slide, lbl, gx + 0.18, gy + 0.1, 5.7, 0.3, size=8, bold=True, color=C_BLUE)
            add_text(slide, str(val), gx + 0.18, gy + 0.45, 5.7, 0.85, size=10, color=C_BG, wrap=True)

        sy = 4.55
        for lbl, key in [("STRATEGY", "strategy"), ("RECOMMENDATIONS", "recommendations")]:
            add_rect(slide, 0.4, sy, 12.5, 0.9, RGBColor(0x0E, 0x16, 0x2E))
            add_text(slide, lbl, 0.55, sy + 0.08, 2.5, 0.3, size=8, bold=True, color=C_BLUE)
            add_text(slide, summary.get(key, "—"), 0.55, sy + 0.38, 12.1, 0.45,
                     size=10, color=C_BG, wrap=True)
            sy += 1.05

        tags = summary.get("tags", [])
        if tags:
            add_text(slide, "  ".join(f"#{t}" for t in tags),
                     0.4, 6.85, 12.5, 0.35, size=10, color=C_BLUE)

    # ── 6. 기회 영역 & 전략 제안 슬라이드 ──
    if summary and summary.get("our_opportunity"):
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, C_BG)
        add_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
        add_text(slide, "05  기회 영역 & 전략 제안", 0.4, 0.2, 12, 0.6, size=18, bold=True, color=C_BG)
        add_text(slide, "경쟁사 약점 기반 우리의 차별화 기회", 0.4, 0.72, 12, 0.32, size=10, color=C_MUTED)

        # 경쟁사 약점 (좌)
        add_rect(slide, 0.3, 1.25, 6.0, 5.8, C_BG2)
        add_rect(slide, 0.3, 1.25, 0.08, 5.8, RGBColor(0xE8, 0x28, 0x4A))
        add_text(slide, "경쟁사 공통 약점", 0.55, 1.35, 5.6, 0.4, size=12, bold=True,
                 color=RGBColor(0xE8, 0x28, 0x4A))
        add_text(slide, summary.get("common_weakness", "—"),
                 0.55, 1.85, 5.6, 2.0, size=11, color=C_DARK, wrap=True)

        # 점수 최하위 항목 추출
        all_scores = {}
        for a in ai_items:
            ce = (a["ai"].get("conversion_elements") or {})
            for k in ["price_emphasis", "product_visibility", "readability",
                      "discount_visibility", "visual_clarity"]:
                v = ce.get(k)
                if v:
                    all_scores.setdefault(k, []).append(int(v))
        score_avgs = {k: round(sum(v)/len(v), 1) for k, v in all_scores.items() if v}
        label_map = {"price_emphasis": "가격 강조", "product_visibility": "제품 가시성",
                     "readability": "가독성", "discount_visibility": "할인 가시성",
                     "visual_clarity": "시각 명확도"}
        if score_avgs:
            add_text(slide, "평균 취약 영역 (1-5점 기준)", 0.55, 3.6, 5.6, 0.35,
                     size=9, bold=True, color=C_MUTED)
            y_sc2 = 4.0
            for k, v in sorted(score_avgs.items(), key=lambda x: x[1])[:4]:
                bar_w = (v / 5) * 4.5
                add_rect(slide, 2.0, y_sc2 + 0.03, bar_w, 0.28,
                         RGBColor(0xE8, 0x28, 0x4A) if v < 3 else C_ORANGE if v < 4 else C_GREEN)
                add_text(slide, label_map.get(k, k), 0.55, y_sc2, 1.4, 0.32, size=9, color=C_DARK)
                add_text(slide, str(v), 2.0 + bar_w + 0.1, y_sc2, 0.4, 0.32, size=9,
                         bold=True, color=C_DARK)
                y_sc2 += 0.42

        # 우리의 기회 (우)
        add_rect(slide, 6.7, 1.25, 6.3, 2.7, C_BG2)
        add_rect(slide, 6.7, 1.25, 0.08, 2.7, C_GREEN)
        add_text(slide, "우리의 기회 영역", 6.95, 1.35, 5.9, 0.4, size=12, bold=True, color=C_GREEN)
        add_text(slide, summary.get("our_opportunity", "—"),
                 6.95, 1.85, 5.9, 2.0, size=11, color=C_DARK, wrap=True)

        add_rect(slide, 6.7, 4.1, 6.3, 2.95, C_BG2)
        add_rect(slide, 6.7, 4.1, 0.08, 2.95, C_BLUE)
        add_text(slide, "소재 제작 제안", 6.95, 4.2, 5.9, 0.4, size=12, bold=True, color=C_BLUE)
        add_text(slide, summary.get("recommendations", "—"),
                 6.95, 4.7, 5.9, 2.2, size=11, color=C_DARK, wrap=True)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ============================================================
# 비교 테이블 뷰
# ============================================================
def render_comparison_table(all_items):
    try:
        import pandas as pd
    except ImportError:
        st.error("pandas가 설치되어 있지 않습니다. requirements.txt에 pandas를 추가하세요.")
        return

    analyzed = [a for a in all_items if a.get("ai") and not (a.get("ai") or {}).get("_error")]
    if not analyzed:
        st.markdown(
            '<div class="empty">'
            '<div class="empty-t">비교할 AI 분석 결과가 없습니다</div>'
            '<div class="empty-d">소재를 선택하고 [선택 분석]을 실행하면<br>여기에 비교 테이블이 생성됩니다.</div>'
            '</div>',
            unsafe_allow_html=True)
        return

    rows = []
    for item in analyzed:
        ai  = item.get("ai") or {}
        ce  = ai.get("conversion_elements") or {}
        vf  = ai.get("visual_facts") or {}
        ma  = ai.get("marketing_analysis") or {}
        cd  = ai.get("creative_diagnosis") or {}
        rows.append({
            "키워드":       item["keyword"],
            "타입":         "이미지" if item["asset_type"] == "image" else "비디오",
            "레이아웃":     (vf.get("layout_type") or ai.get("layout_type") or "—")[:20],
            "후크":         (ai.get("hook") or ma.get("hook_type") or "—")[:20],
            "소구":         str(ai.get("appeal") or "—")[:30],
            "타겟":         (ai.get("target") or "—")[:40],
            "핵심메시지":   (ai.get("message") or "—")[:60],
            "가격강조":     int(ce.get("price_emphasis") or 0),
            "제품가시성":   int(ce.get("product_visibility") or 0),
            "가독성":       int(ce.get("readability") or 0),
            "할인가시성":   int(ce.get("discount_visibility") or 0),
            "시각명확도":   int(ce.get("visual_clarity") or 0),
            "전환력":       int(ce.get("overall_conversion_power") or 0),
            "개선방향":     (cd.get("improvement_direction") or "—")[:80],
            "썸네일":       item["image_url"],
        })

    df = pd.DataFrame(rows)
    df = df.sort_values("전환력", ascending=False).reset_index(drop=True)

    st.markdown(
        f'<div class="sec"><div class="sec-t">소재 비교 테이블</div>'
        f'<div class="sec-n">{len(analyzed)}개 · AI 분석 완료 소재 · 전환력 기준 정렬</div></div>',
        unsafe_allow_html=True)

    st.dataframe(
        df,
        column_config={
            "썸네일":     st.column_config.ImageColumn("썸네일", width="small"),
            "가격강조":   st.column_config.NumberColumn("가격강조",   format="%d ⭐", min_value=0, max_value=5),
            "제품가시성": st.column_config.NumberColumn("제품가시성", format="%d ⭐", min_value=0, max_value=5),
            "가독성":     st.column_config.NumberColumn("가독성",     format="%d ⭐", min_value=0, max_value=5),
            "할인가시성": st.column_config.NumberColumn("할인가시성", format="%d ⭐", min_value=0, max_value=5),
            "시각명확도": st.column_config.NumberColumn("시각명확도", format="%d ⭐", min_value=0, max_value=5),
            "전환력":     st.column_config.ProgressColumn("전환력", min_value=0, max_value=5, format="%d"),
            "핵심메시지": st.column_config.TextColumn("핵심메시지", width="large"),
            "개선방향":   st.column_config.TextColumn("개선방향",   width="large"),
        },
        use_container_width=True,
        height=420,
    )

    # 분포 차트
    st.markdown("<div style='height:20px'></div>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns(3)
    with c1:
        st.markdown('<div class="slbl">소구 유형 분포</div>', unsafe_allow_html=True)
        appeal_counts = df["소구"].value_counts()
        st.bar_chart(appeal_counts, height=220)
    with c2:
        st.markdown('<div class="slbl">후크 유형 분포</div>', unsafe_allow_html=True)
        hook_counts = df["후크"].value_counts()
        st.bar_chart(hook_counts, height=220)
    with c3:
        st.markdown('<div class="slbl">전환력 점수 분포</div>', unsafe_allow_html=True)
        score_dist = df["전환력"].value_counts().sort_index()
        st.bar_chart(score_dist, height=220)

    # 평균 점수 요약
    st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)
    score_cols = ["가격강조", "제품가시성", "가독성", "할인가시성", "시각명확도", "전환력"]
    avg_scores = df[score_cols].mean().round(2)
    cols_s = st.columns(len(score_cols))
    for col, (label, val) in zip(cols_s, avg_scores.items()):
        color = "var(--ok)" if val >= 4 else "var(--wn)" if val >= 3 else "var(--er)"
        col.markdown(
            f'<div class="sc" style="padding:14px 16px;">'
            f'<div class="sc-lbl">AVG {label.upper()}</div>'
            f'<div class="sc-val" style="font-size:24px;color:{color}">{val}</div>'
            f'<div class="sc-sub">/ 5점</div></div>',
            unsafe_allow_html=True)


# ============================================================
# 헤더
# ============================================================
st.markdown("""
<div class="hdr">
  <div class="hdr-row">
    <div class="logo">AD<b>INTEL</b></div>
    <div class="ver">v4.0</div>
    <span class="persist-badge">💾 DB 자동저장</span>
  </div>
  <div class="sub">META AD LIBRARY INTELLIGENCE BOARD</div>
</div>""", unsafe_allow_html=True)


# ============================================================
# 사이드바
# ============================================================
with st.sidebar:
    st.markdown('<div class="slbl">DISPLAY</div>', unsafe_allow_html=True)
    new_dark = st.toggle("Dark Mode", value=st.session_state.dark)
    if new_dark != st.session_state.dark:
        st.session_state.dark = new_dark
        st.rerun()
    st.markdown("---")

    st.markdown('<div class="slbl">AI ANALYSIS</div>', unsafe_allow_html=True)
    ai_on = st.toggle("소구포인트 분석 ON/OFF", value=st.session_state.ai_on)
    st.session_state.ai_on = ai_on
    if ai_on:
        if API_KEY:
            st.success("API Key 연결됨")

            # 분석 모델 선택
            st.markdown('<div class="slbl">분석 모델 선택</div>', unsafe_allow_html=True)
            model_choice = st.radio(
                "model_radio",
                options=["haiku", "sonnet"],
                format_func=lambda x: (
                    "⚡ Haiku — 빠름·저비용 (~1원)" if x == "haiku"
                    else "🎯 Sonnet — 정밀·고품질 (~15원)"
                ),
                index=0 if st.session_state.get("analysis_model", "haiku") == "haiku" else 1,
                label_visibility="collapsed",
            )
            if model_choice != st.session_state.get("analysis_model"):
                st.session_state.analysis_model = model_choice
            if model_choice == "sonnet":
                st.info("비전 분석 품질이 크게 높습니다.\n소재당 약 10~20원 수준.")

            if st.button("API 연결 테스트", use_container_width=True):
                ok, msg = test_api()
                if ok:
                    st.success(msg)
                else:
                    st.error(msg)
        else:
            st.error("API Key 없음")
            st.caption("Streamlit Cloud > Settings > Secrets\nANTHROPIC_API_KEY = 'sk-ant-...'")
    else:
        st.caption("OFF — API 호출 없음 (무료)")
    st.markdown("---")

    st.markdown('<div class="slbl">SCRAPE</div>', unsafe_allow_html=True)
    country = st.selectbox("국가", ["KR", "US", "JP", "GB", "SG", "AU"], index=0, label_visibility="collapsed")
    st.caption(f"대상 국가: **{country}**")
    scrolls = st.slider("스크롤 깊이", 3, 15, 8, 1)
    max_n   = st.slider("최대 수집량", 10, 150, 60, 10)
    st.markdown("---")

    st.markdown('<div class="slbl">FILTER & SORT</div>', unsafe_allow_html=True)
    ftype = st.selectbox("소재 타입", ["전체", "image", "video_poster"], label_visibility="collapsed")
    kws   = ["전체"] + list(dict.fromkeys(a["keyword"] for a in st.session_state.assets))
    fkw   = st.selectbox("키워드", kws, label_visibility="collapsed")
    fstar = st.toggle("즐겨찾기만", value=False)
    fai   = st.toggle("AI 분석된 것만", value=False)
    cols  = st.select_slider("열 수", options=[2, 3, 4, 5], value=4)
    sort  = st.selectbox("정렬", ["최신순", "오래된순", "키워드순", "즐겨찾기순", "전환력순"], label_visibility="collapsed")
    st.markdown("---")

    st.markdown('<div class="slbl">EXPORT</div>', unsafe_allow_html=True)
    exp_items = get_visible(ftype, fkw if fkw != "전체" else None, fstar, fai)
    stars     = [a for a in st.session_state.assets if a.get("starred")]
    if exp_items:
        st.download_button("CSV 내보내기", data=to_csv(exp_items),
            file_name=f"adintel_{time.strftime('%Y%m%d_%H%M')}.csv",
            mime="text/csv", use_container_width=True)
    if stars:
        st.download_button("즐겨찾기 CSV", data=to_csv(stars),
            file_name=f"adintel_star_{time.strftime('%Y%m%d_%H%M')}.csv",
            mime="text/csv", use_container_width=True)

    st.markdown('<div style="height:4px"></div>', unsafe_allow_html=True)
    if st.button("보드 초기화", use_container_width=True):
        do_reset()
        st.rerun()

    st.markdown("---")
    st.markdown('<div class="slbl">DB 상태</div>', unsafe_allow_html=True)
    db_count = len(st.session_state.assets)
    st.caption(f"저장된 소재: **{db_count}개** | DB: `{DB_PATH}`")
    if st.button("DB 새로고침", use_container_width=True):
        st.session_state.assets = db_load_assets()
        st.session_state.hidden = db_load_hidden()
        st.session_state.history = db_load_history()
        st.session_state.summary = db_load_summary()
        st.rerun()

    if st.session_state.log:
        st.markdown("---")
        st.markdown('<div class="slbl">LOG</div>', unsafe_allow_html=True)
        for lg in reversed(st.session_state.log[-8:]):
            cls = "ok" if lg["ok"] else "er"
            st.markdown(
                f'<div class="log"><span class="ts">{lg["t"]}</span>'
                f'<span class="{cls}">{"+" if lg["ok"] else "!"} {lg["kw"]} — {lg["n"]}건</span></div>',
                unsafe_allow_html=True)


# ============================================================
# AI 배너
# ============================================================
if ai_on and API_KEY:
    st.markdown("""<div class="banner">
        <div class="dot dot-on"></div>
        <span class="banner-txt">AI 분석 <b>ON</b> &nbsp;—&nbsp; 선택 소재만 캡처 분석 · Claude Vision · 분석 결과 DB 자동저장</span>
    </div>""", unsafe_allow_html=True)
elif ai_on and not API_KEY:
    st.markdown("""<div class="banner banner-w">
        <div class="dot" style="background:var(--wn)"></div>
        <span class="banner-txt">AI ON — API Key 없음 &nbsp;·&nbsp; Settings > Secrets > <b>ANTHROPIC_API_KEY</b></span>
    </div>""", unsafe_allow_html=True)


# ============================================================
# 검색창
# ============================================================
st.markdown('<div class="srch">', unsafe_allow_html=True)
st.markdown('<div class="srch-lbl">KEYWORD SEARCH</div>', unsafe_allow_html=True)
c1, c2, c3 = st.columns([5, 1, 1])
with c1:
    kw = st.text_input("kw",
        placeholder="예: 캐리어, 공기청정기, 스킨케어, 다이어트",
        label_visibility="collapsed",
        on_change=lambda: st.session_state.update({"_enter": True}))
with c2:
    do_new = st.button("검색", use_container_width=True)
with c3:
    do_add = st.button("누적 검색", use_container_width=True)

if st.session_state.pop("_enter", False):
    do_new = True
st.caption("엔터 / 검색 = 새 검색   ·   누적 검색 = 기존 보드에 추가   ·   💾 수집 결과 자동 DB저장")
st.markdown('</div>', unsafe_allow_html=True)


# ============================================================
# 검색 실행
# ============================================================
if do_new or do_add:
    q = (kw or "").strip()
    if not q:
        st.warning("검색어를 입력하세요.")
    else:
        try:
            if do_new:
                st.session_state.assets   = []
                st.session_state.hidden   = set()
                st.session_state.history  = []
                st.session_state.summary  = None
                st.session_state.selected = set()

            pb = st.progress(0, text=f"'{q}' — Meta Ad Library 접근 중...")
            items = scrape(q, country, scrolls, max_n)
            pb.progress(0.7, text=f"수집 완료 {len(items)}개 · DB 저장 중...")

            merged, added = merge(st.session_state.assets, items)
            st.session_state.assets = merged

            # 신규 소재 DB 저장
            new_only = [a for a in merged if any(n["id"] == a["id"] for n in items)]
            db_upsert_assets(new_only)

            if q not in st.session_state.history:
                st.session_state.history.append(q)
                db_add_history(q)

            pb.progress(1.0, text="완료!")
            pb.empty()

            st.session_state.log.append({"t": time.strftime("%H:%M"), "kw": q, "n": added, "ok": True})
            st.success(f"{added}개 수집 · DB 저장 완료") if added > 0 else st.info("새 소재 없음")
            st.rerun()

        except Exception as e:
            st.session_state.log.append({"t": time.strftime("%H:%M"), "kw": q, "n": 0, "ok": False})
            st.error(f"오류: {e}")


# ============================================================
# 스탯
# ============================================================
all_a = st.session_state.assets
shown = get_visible(ftype, fkw if fkw != "전체" else None, fstar, fai)

if sort == "최신순":       shown = sorted(shown, key=lambda a: a["created_at"], reverse=True)
elif sort == "오래된순":   shown = sorted(shown, key=lambda a: a["created_at"])
elif sort == "키워드순":   shown = sorted(shown, key=lambda a: a["keyword"])
elif sort == "즐겨찾기순": shown = sorted(shown, key=lambda a: (not a.get("starred"), a["created_at"]))
elif sort == "전환력순":
    def conv_score(a):
        ai = a.get("ai") or {}
        ce = ai.get("conversion_elements") or {}
        return -(ce.get("overall_conversion_power") or 0)
    shown = sorted(shown, key=conv_score)

ic = sum(1 for a in shown if a["asset_type"] == "image")
vc = sum(1 for a in shown if a["asset_type"] == "video_poster")

for col, (lbl, val, sub) in zip(st.columns(6), [
    ("TOTAL",    len(all_a), "누적 수집"),
    ("SHOWN",    len(shown), f"IMG {ic} · VID {vc}"),
    ("KEYWORDS", len(set(a["keyword"] for a in all_a)), "검색어"),
    ("STARRED",  sum(1 for a in all_a if a.get("starred")), "즐겨찾기"),
    ("HIDDEN",   len(st.session_state.hidden), "제외"),
    ("AI",       sum(1 for a in all_a if a.get("ai")), "분석 완료"),
]):
    col.markdown(
        f'<div class="sc"><div class="sc-lbl">{lbl}</div>'
        f'<div class="sc-val">{val}</div><div class="sc-sub">{sub}</div></div>',
        unsafe_allow_html=True)

st.markdown('<div style="height:24px"></div>', unsafe_allow_html=True)


# ============================================================
# 키워드 태그
# ============================================================
if st.session_state.history:
    html = '<div class="tags">'
    for h in st.session_state.history[-20:]:
        n = sum(1 for a in all_a if a["keyword"] == h)
        html += f'<span class="tag"><span class="tag-dot"></span>{h} <span style="opacity:.5">({n})</span></span>'
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)


# ============================================================
# 총합 인사이트 패널
# ============================================================
summary = st.session_state.get("summary")
if summary:
    tags_html = "".join(f'<span class="summary-tag">{t}</span>' for t in summary.get("tags", []))
    st.markdown(
        '<div class="summary-box">'
        '<div class="summary-head">TOTAL INSIGHT — 선택 소재 종합 분석</div>'
        '<div class="summary-grid">'
        '<div class="summary-item"><div class="summary-item-lbl">DOMINANT APPEAL</div>'
        f'<div class="summary-item-val">{summary.get("dominant_appeal","—")}</div></div>'
        '<div class="summary-item"><div class="summary-item-lbl">COMMON TARGET</div>'
        f'<div class="summary-item-val">{summary.get("common_target","—")}</div></div>'
        '<div class="summary-item"><div class="summary-item-lbl">KEY MESSAGE PATTERN</div>'
        f'<div class="summary-item-val">{summary.get("key_message","—")}</div></div>'
        '<div class="summary-item"><div class="summary-item-lbl">COMMON WEAKNESS</div>'
        f'<div class="summary-item-val">{summary.get("common_weakness","—")}</div></div>'
        '</div>'
        '<div class="summary-strategy"><div class="summary-strategy-lbl">OUR OPPORTUNITY</div>'
        f'<div class="summary-strategy-val">{summary.get("our_opportunity","—")}</div></div>'
        '<div class="summary-strategy"><div class="summary-strategy-lbl">STRATEGY</div>'
        f'<div class="summary-strategy-val">{summary.get("strategy","—")}</div></div>'
        '<div class="summary-strategy" style="margin-bottom:0">'
        '<div class="summary-strategy-lbl">RECOMMENDATIONS</div>'
        f'<div class="summary-strategy-val">{summary.get("recommendations","—")}</div></div>'
        f'<div style="margin-top:12px"><div class="summary-tags">{tags_html}</div></div>'
        '</div>',
        unsafe_allow_html=True)
    if st.button("인사이트 초기화", key="clear_summary"):
        st.session_state.summary  = None
        st.session_state.selected = set()
        st.rerun()


# ============================================================
# 메인 탭 — 소재 보드 / 비교 테이블
# ============================================================
tab_board, tab_table = st.tabs(["◼ 소재 보드", "⊞ 비교 테이블"])


# ── 탭1: 소재 보드 ──────────────────────────────────────────
with tab_board:
    sel = st.session_state.selected

    st.markdown(
        f'<div class="sec"><div class="sec-t">소재 보드</div>'
        f'<div class="sec-n">{len(shown)} assets · {sort}</div></div>',
        unsafe_allow_html=True)

    if ai_on and API_KEY and shown:
        sel_items      = [a for a in st.session_state.assets if a["id"] in sel]
        sel_unanalyzed = [a for a in sel_items if not a.get("ai")]
        analyzed_sel   = [a for a in sel_items if a.get("ai") and not (a.get("ai") or {}).get("_error")]

        bar_c, btn_c1, btn_c2, btn_c3 = st.columns([3, 1, 1, 1])
        with bar_c:
            st.markdown(
                f'<div class="sel-bar"><span class="sel-count">{len(sel)}</span>'
                '<span class="sel-bar-txt">선택됨  ·  ○ 버튼으로 소재를 선택하세요</span></div>',
                unsafe_allow_html=True)
        with btn_c1:
            btn_label = f"선택 분석 ({len(sel_unanalyzed)})" if sel_unanalyzed else f"선택 재분석 ({len(sel_items)})"
            do_analyze = st.button(btn_label, use_container_width=True, disabled=not sel_items)
        with btn_c2:
            do_insight = (
                st.button(f"종합 인사이트 ({len(analyzed_sel)})", use_container_width=True)
                if analyzed_sel else False
            )
        with btn_c3:
            do_ppt = (
                st.button(f"PPT 내보내기 ({len(analyzed_sel)})", use_container_width=True)
                if analyzed_sel else False
            )

        # ── 분석 애니메이션: 컬럼 밖 전체 너비 ──
        status = st.empty()

        if do_analyze and sel_items:
            target_items = sel_unanalyzed if sel_unanalyzed else sel_items
            for it in target_items:
                it["ai"] = None
                it["img_b64"] = ""
            n = len(target_items)

            status.markdown(
                _analysis_status_html("capture", 0, n, "Playwright로 광고 카드 전체 영역을 캡처합니다"),
                unsafe_allow_html=True,
            )
            capture_screenshots_for_items(target_items, scrolls=scrolls)
            captured = sum(1 for it in target_items if it.get("img_b64"))

            status.markdown(
                _analysis_status_html(
                    "analyze", captured, n,
                    f"캡처 완료 {captured}개 · OCR + Vision 분석 병렬 실행 중",
                ),
                unsafe_allow_html=True,
            )
            _, errors = analyze_parallel(target_items, max_workers=3, force=True)
            done = sum(1 for it in target_items if it.get("ai") and not (it.get("ai") or {}).get("_error"))

            status.markdown(
                _analysis_status_html("done", done, n),
                unsafe_allow_html=True,
            )
            time.sleep(1.8)
            status.empty()

            if errors:
                with st.expander(f"⚠ 분석 오류 {len(errors)}건", expanded=False):
                    for e in errors:
                        st.markdown(f'<div class="err-box">{e}</div>', unsafe_allow_html=True)

            id_map = {it["id"]: it for it in target_items}
            for a in st.session_state.assets:
                if a["id"] in id_map:
                    a["ai"] = id_map[a["id"]].get("ai")
            st.rerun()

        if do_insight:
            with st.spinner("종합 인사이트 생성 중..."):
                summ = summarize_insights(analyzed_sel)
                if summ:
                    st.session_state.summary = summ
                    db_save_summary(summ)
            st.rerun()

        if do_ppt:
            with st.spinner(f"PPT 생성 중 ({len(analyzed_sel)}개 소재)..."):
                pptx_bytes = to_pptx(analyzed_sel, st.session_state.get("summary"))
            fname = f"adintel_{time.strftime('%Y%m%d_%H%M')}.pptx"
            st.download_button(
                "다운로드",
                data=pptx_bytes,
                file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

    if not shown:
        st.markdown(
            '<div class="empty"><div class="empty-t">수집된 소재가 없습니다</div>'
            '<div class="empty-d">키워드를 입력하고 검색하면<br>'
            'Meta Ad Library에서 광고 소재를 실시간으로 수집합니다.</div></div>',
            unsafe_allow_html=True)
    else:
        grid = st.columns(cols)
        for i, item in enumerate(shown):
            with grid[i % cols]:
                is_sel = item["id"] in sel

                st.markdown('<div class="card">', unsafe_allow_html=True)
                try:
                    st.image(item["image_url"], use_container_width=True)
                except Exception:
                    st.markdown(
                        '<div style="height:100px;background:var(--bg3);display:flex;'
                        'align-items:center;justify-content:center;color:var(--mu);font-size:10px;">'
                        'LOAD FAILED</div>', unsafe_allow_html=True)

                bc   = "b-img" if item["asset_type"] == "image" else "b-vid"
                bl   = "IMG"   if item["asset_type"] == "image" else "VID"
                sel_b = '<span class="bdg" style="background:var(--ac);color:#fff;border-color:var(--ac)">SEL</span> ' if is_sel else ""
                sv_b  = '<span class="bdg b-sav">★</span> ' if item.get("starred") else ""
                ai_b  = '<span class="bdg b-ai">AI</span> ' if (item.get("ai") and not (item.get("ai") or {}).get("_error")) else ""
                er_b  = '<span class="bdg" style="background:rgba(232,40,74,.1);color:var(--er);border-color:rgba(232,40,74,.2)">ERR</span> ' if (item.get("ai") and (item.get("ai") or {}).get("_error")) else ""
                cap   = item.get("caption", "")
                cap_h = f'<div class="card-cap">"{cap[:55]}{"..." if len(cap) > 55 else ""}"</div>' if cap else ""

                ai_data = _sanitize_json_deep(item.get("ai")) if item.get("ai") else None
                conv_score_str = ""
                if ai_data and not ai_data.get("_error"):
                    ce = ai_data.get("conversion_elements") or {}
                    s = ce.get("overall_conversion_power")
                    if s:
                        conv_score_str = f' · ⚡{s}/5'

                st.markdown(
                    '<div class="card-body">'
                    f'<div class="card-kw">{item["keyword"]} · {item["country"]}{conv_score_str}</div>'
                    + sel_b + sv_b + ai_b + er_b +
                    f'<span class="bdg {bc}">{bl}</span>'
                    + cap_h +
                    f'<div class="card-meta">{item["width"]}x{item["height"]}px · {item["created_at"][11:16]}</div>'
                    '</div>',
                    unsafe_allow_html=True)

                if ai_data:
                    if ai_data.get("_error"):
                        st.markdown(
                            '<div class="ai-wrap"><div class="ai-box" style="border-color:var(--er)">'
                            f'<div class="ai-head" style="color:var(--er)">분석 오류</div>'
                            f'<div class="ai-body" style="font-size:11px">{ai_data["_error"]}</div>'
                            '</div></div>',
                            unsafe_allow_html=True)
                    else:
                        vf = ai_data.get("visual_facts") or {}
                        ma = ai_data.get("marketing_analysis") or {}
                        ce = ai_data.get("conversion_elements") or {}
                        cd = ai_data.get("creative_diagnosis") or {}
                        has_img = bool(db_get_img_b64(item["id"]))

                        # 새 구조 (observe/interpret) 우선, 구버전 호환 fallback
                        ob = ai_data.get("observe") or {}
                        ip = ai_data.get("interpret") or {}

                        visible_text = _safe_join(ob.get("text_exact", []), " / ") or _safe_join(vf.get("visible_text", []), " / ") or ai_data.get("copy", "—")
                        main_visual  = ob.get("main_visual_description") or ai_data.get("main_visual", "—")
                        layout       = ob.get("layout_type") or ai_data.get("layout_type", "—")
                        colors_str   = _safe_join(ob.get("color_palette", [])) or _safe_join(vf.get("colors", [])) or "—"
                        hook         = ip.get("hook_type") or ai_data.get("hook", "—")
                        hook_ev      = ip.get("hook_evidence", "")
                        appeal       = ip.get("primary_appeal") or ai_data.get("appeal", "—")
                        appeal_ev    = ip.get("appeal_evidence", "")
                        target       = ip.get("target_audience") or ai_data.get("target", "—")
                        target_ev    = ip.get("target_evidence", "")
                        message      = ip.get("core_message") or ai_data.get("message", "—")
                        strengths    = _safe_join((cd.get("strengths") or []))  or "—"
                        weaknesses   = _safe_join((cd.get("weaknesses") or [])) or "—"
                        improve      = cd.get("improvement_direction") or dg.get("improvement_direction") if (dg := ai_data.get("diagnosis") or {}) else cd.get("improvement_direction") or "—"

                        score_line = (
                            f'가격 {ce.get("price_emphasis","—")} / '
                            f'제품 {ce.get("product_visibility","—")} / '
                            f'가독성 {ce.get("readability","—")} / '
                            f'전환력 {ce.get("overall_conversion_power","—")}'
                        )
                        cap_src = item.get("capture_source", "")
                        capture_badge = (
                            '<span class="bdg b-ai" style="font-size:8px">카드캡처</span> ' if "card" in cap_src
                            else '<span class="bdg b-vid" style="font-size:8px">이미지캡처</span> ' if has_img
                            else '<span class="bdg" style="background:var(--bg3);color:var(--mu);font-size:8px">캡처실패</span> '
                        )
                        tags_str = "".join(f'<span class="ai-tag">{t}</span>' for t in ai_data.get("tags", []))

                        def ev_span(ev):
                            if not ev or ev == "확인 불가":
                                return ""
                            return f'<span style="font-size:10px;color:var(--mu);display:block;padding-left:4px;border-left:2px solid var(--bd2);margin:1px 0 4px 0">근거: {ev[:80]}</span>'

                        st.markdown(
                            '<div class="ai-wrap"><div class="ai-box">'
                            f'<div class="ai-head">OBSERVE → INTERPRET {capture_badge}</div>'
                            '<div class="ai-body">'
                            '<span style="font-size:9px;color:var(--ac);letter-spacing:1px">OBSERVE</span><br>'
                            f'<b>실제문구</b>&nbsp;{visible_text}<br>'
                            f'<b>비주얼</b>&nbsp;&nbsp;&nbsp;{main_visual}<br>'
                            f'<b>레이아웃</b>&nbsp;{layout}<br>'
                            f'<b>색감</b>&nbsp;&nbsp;&nbsp;&nbsp;{colors_str}<br>'
                            '<span style="font-size:9px;color:var(--ac);letter-spacing:1px;margin-top:6px;display:block">INTERPRET (관찰 근거 기반)</span>'
                            f'<b>후크</b>&nbsp;&nbsp;&nbsp;{hook}<br>{ev_span(hook_ev)}'
                            f'<b>소구</b>&nbsp;&nbsp;&nbsp;{appeal}<br>{ev_span(appeal_ev)}'
                            f'<b>타겟</b>&nbsp;&nbsp;&nbsp;{target}<br>{ev_span(target_ev)}'
                            f'<b>메시지</b>&nbsp;{message}<br>'
                            f'<b>점수</b>&nbsp;&nbsp;&nbsp;{score_line}<br>'
                            f'<b>장점</b>&nbsp;&nbsp;&nbsp;{strengths}<br>'
                            f'<b>약점</b>&nbsp;&nbsp;&nbsp;{weaknesses}<br>'
                            f'<b>개선</b>&nbsp;&nbsp;&nbsp;{improve}<br>'
                            f'<div style="margin-top:6px">{tags_str}</div>'
                            '</div></div></div>',
                            unsafe_allow_html=True)

                st.markdown('</div>', unsafe_allow_html=True)

                b1, b2, b3, b4 = st.columns(4)
                with b1:
                    if st.button("✓" if is_sel else "○", key="sel_" + item["id"], use_container_width=True):
                        if is_sel:
                            st.session_state.selected.discard(item["id"])
                        else:
                            st.session_state.selected.add(item["id"])
                        st.rerun()
                with b2:
                    if st.button("★" if item.get("starred") else "☆", key="sv_" + item["id"], use_container_width=True):
                        toggle_star(item["id"])
                        st.rerun()
                with b3:
                    if st.button("✕", key="hd_" + item["id"], use_container_width=True):
                        st.session_state.hidden.add(item["id"])
                        db_add_hidden(item["id"])
                        st.rerun()
                with b4:
                    st.link_button("↗", item["source_url"], use_container_width=True)


# ── 탭2: 비교 테이블 ─────────────────────────────────────────
with tab_table:
    render_comparison_table(all_a)


# ============================================================
# 푸터
# ============================================================
st.markdown("<div style='height:40px'></div>", unsafe_allow_html=True)
th   = "DARK" if D else "LIGHT"
model_label = "Sonnet" if st.session_state.get("analysis_model")=="sonnet" else "Haiku"
ai_s = f"AI ON · {model_label}" if (ai_on and API_KEY) else "AI OFF"
st.markdown(
    f'<div style="border-top:2px solid var(--bd);padding-top:14px;display:flex;justify-content:space-between;">'
    f'<span style="font-size:10px;color:var(--mu)">ADINTEL v4.0 · {th} · {ai_s} · SQLite 자동저장</span>'
    f'<span style="font-size:10px;color:var(--mu)">{time.strftime("%Y-%m-%d")}</span>'
    f'</div>',
    unsafe_allow_html=True)
